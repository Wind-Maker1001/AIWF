from __future__ import annotations

import os
from typing import Any, Dict, List, Optional, Tuple

from aiwf.office_style import (
    office_lang,
    office_text,
    office_font_name,
    office_theme_settings,
    office_layout_settings,
    office_quality_mode,
    office_is_high_quality,
    hex_to_rgb,
    pil_font,
    docx_apply_font,
    add_picture_fit,
)


def write_profile_illustration_png(
    path: str,
    profile: Dict[str, Any],
    params: Optional[Dict[str, Any]] = None,
    *,
    utc_now_str,
) -> bool:
    try:
        from PIL import Image, ImageDraw  # type: ignore
    except Exception:
        return False

    image_params = dict(params or {})
    if office_lang(params) == "zh" and pil_font(20, params) is None:
        image_params["office_lang"] = "en"
    theme = office_theme_settings(image_params)
    width = 1200
    height = 720
    img = Image.new("RGB", (width, height), color=theme["bg_rgb"])
    draw = ImageDraw.Draw(img)
    font_title = pil_font(40, image_params, bold=True)
    font_time = pil_font(20, image_params)
    font_card_label = pil_font(24, image_params)
    font_card_value = pil_font(34, image_params, bold=True)
    font_chart = pil_font(22, image_params)
    font_bar_label = pil_font(18, image_params)
    font_bar_value = pil_font(20, image_params, bold=True)

    head_hex = str(theme.get("primary_hex", "1F4E78"))
    head_rgb = tuple(int(head_hex[i : i + 2], 16) for i in (0, 2, 4))
    draw.rectangle([(0, 0), (width, 88)], fill=head_rgb)
    draw.text(
        (32, 28),
        f"{theme.get('report_title')} - {office_text('质量快照', 'Snapshot', image_params)}",
        fill=(255, 255, 255),
        font=font_title,
    )
    draw.text((820, 32), utc_now_str(), fill=(220, 230, 240), font=font_time)

    quality = profile.get("quality") if isinstance(profile.get("quality"), dict) else {}
    kpis = [
        (office_text("总行数", "Rows", image_params), int(profile.get("rows", 0) or 0)),
        (office_text("输入行", "Input", image_params), int(quality.get("input_rows", 0) or 0)),
        (office_text("输出行", "Output", image_params), int(quality.get("output_rows", 0) or 0)),
        (office_text("无效行", "Invalid", image_params), int(quality.get("invalid_rows", 0) or 0)),
        (office_text("过滤行", "Filtered", image_params), int(quality.get("filtered_rows", 0) or 0)),
    ]

    card_w = 210
    card_h = 116
    gap = 18
    start_x = 36
    y = 118
    for i, (name, val) in enumerate(kpis):
        x = start_x + i * (card_w + gap)
        draw.rounded_rectangle([(x, y), (x + card_w, y + card_h)], radius=12, fill=(255, 255, 255), outline=(210, 218, 230))
        draw.text((x + 16, y + 18), name, fill=(80, 90, 110), font=font_card_label)
        draw.text((x + 16, y + 58), str(val), fill=(26, 34, 50), font=font_card_value)

    chart_x = 72
    chart_y = 320
    chart_w = 1050
    chart_h = 300
    draw.rectangle([(chart_x, chart_y), (chart_x + chart_w, chart_y + chart_h)], fill=(255, 255, 255), outline=(210, 218, 230))
    draw.text(
        (chart_x + 18, chart_y + 12),
        office_text("流程数据量", "Pipeline Volume", image_params),
        fill=(55, 65, 85),
        font=font_chart,
    )

    bar_vals = [
        (office_text("输入", "Input", image_params), int(quality.get("input_rows", 0) or 0), (70, 130, 180)),
        (office_text("输出", "Output", image_params), int(quality.get("output_rows", 0) or 0), (46, 164, 79)),
        (office_text("无效", "Invalid", image_params), int(quality.get("invalid_rows", 0) or 0), (215, 70, 70)),
        (office_text("过滤", "Filtered", image_params), int(quality.get("filtered_rows", 0) or 0), (235, 160, 60)),
    ]
    max_v = max([v for _, v, _ in bar_vals] + [1])
    usable_h = chart_h - 90
    bar_w = 160
    bar_gap = 68
    bx = chart_x + 80
    base_y = chart_y + chart_h - 38
    for name, val, color in bar_vals:
        h = int((float(val) / float(max_v)) * usable_h)
        draw.rectangle([(bx, base_y - h), (bx + bar_w, base_y)], fill=color)
        draw.text((bx + 20, base_y + 8), name, fill=(70, 80, 100), font=font_bar_label)
        draw.text((bx + 20, base_y - h - 24), str(val), fill=(40, 50, 70), font=font_bar_value)
        bx += bar_w + bar_gap

    os.makedirs(os.path.dirname(path), exist_ok=True)
    img.save(path)
    return True


def write_fin_xlsx(
    xlsx_path: str,
    rows: List[Dict[str, Any]],
    image_path: Optional[str] = None,
    params: Optional[Dict[str, Any]] = None,
    *,
    to_decimal,
    build_profile,
    utc_now_str,
) -> None:
    from openpyxl import Workbook  # type: ignore
    from openpyxl.styles import Alignment, Font, PatternFill  # type: ignore
    from openpyxl.utils import get_column_letter  # type: ignore

    wb = Workbook()
    ws = wb.active
    ws.title = "detail"
    if rows:
        columns = list(rows[0].keys())
        seen = set(columns)
        for r in rows[1:]:
            for k in r.keys():
                if k not in seen:
                    columns.append(k)
                    seen.add(k)
    else:
        columns = ["id", "amount"]

    theme = office_theme_settings(params)
    layout = office_layout_settings(params)
    high_quality = office_is_high_quality(params)
    font_name = office_font_name(params)
    header_fill = PatternFill(fill_type="solid", fgColor=str(theme.get("primary_hex", "1F4E78")))
    alt_fill = PatternFill(fill_type="solid", fgColor="F4F7FB")
    header_font = Font(name=font_name, color="FFFFFF", bold=True)
    body_font = Font(name=font_name)
    title_font = Font(name=font_name, bold=True, color=str(theme.get("primary_hex", "1F4E78")), size=14)
    center = Alignment(horizontal="center", vertical="center")
    left = Alignment(horizontal="left", vertical="center")

    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=max(2, len(columns)))
    ws.cell(row=1, column=1, value=str(theme.get("report_title"))).font = title_font
    ws.cell(row=1, column=1).alignment = left
    ws.row_dimensions[1].height = 24

    header_row = 2
    for col_idx, c in enumerate(columns, start=1):
        cell = ws.cell(row=header_row, column=col_idx, value=c)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = center
    for row_idx, r in enumerate(rows, start=header_row + 1):
        for col_idx, c in enumerate(columns, start=1):
            cell = ws.cell(row=row_idx, column=col_idx, value=r.get(c))
            cell.font = body_font
            cell.alignment = left
            if high_quality and row_idx % 2 == 1:
                cell.fill = alt_fill

    ws.freeze_panes = "A3"
    if columns:
        ws.auto_filter.ref = (
            f"A{header_row}:{get_column_letter(len(columns))}{max(header_row + 1, len(rows) + header_row)}"
        )

    numeric_fields = []
    for c in columns:
        if c.lower() in {"amount", "sum", "min", "max", "avg"}:
            numeric_fields.append(c)
            continue
        for r in rows[:100]:
            if to_decimal(r.get(c)) is not None:
                numeric_fields.append(c)
                break
    for c in numeric_fields:
        col_idx = columns.index(c) + 1
        for row_idx in range(header_row + 1, len(rows) + header_row + 1):
            ws.cell(row=row_idx, column=col_idx).number_format = "#,##0.00"

    for col_idx, c in enumerate(columns, start=1):
        max_len = len(str(c))
        for row in rows[:500]:
            max_len = max(max_len, len(str(row.get(c) if row.get(c) is not None else "")))
        ws.column_dimensions[get_column_letter(col_idx)].width = min(max_len + 2, 48)

    profile_like = build_profile(rows, {"input_rows": len(rows), "output_rows": len(rows)}, "xlsx.export")
    summary = wb.create_sheet("summary")
    summary["A1"] = office_text("指标", "Metric", params)
    summary["B1"] = office_text("数值", "Value", params)
    summary["A1"].fill = header_fill
    summary["B1"].fill = header_fill
    summary["A1"].font = header_font
    summary["B1"].font = header_font
    summary["A1"].alignment = center
    summary["B1"].alignment = center
    metrics = [
        (office_text("主题", "Theme", params), theme.get("display_name", theme.get("name"))),
        (office_text("质量模式", "Quality Mode", params), office_quality_mode(params)),
        (office_text("行数", "Rows", params), profile_like.get("rows")),
        (office_text("列数", "Columns", params), len(columns)),
        (office_text("金额总计", "Sum Amount", params), profile_like.get("sum_amount")),
        (office_text("金额最小值", "Min Amount", params), profile_like.get("min_amount")),
        (office_text("金额最大值", "Max Amount", params), profile_like.get("max_amount")),
        (office_text("金额均值", "Avg Amount", params), profile_like.get("avg_amount")),
        (office_text("生成时间", "Generated At", params), utc_now_str()),
    ]
    for i, (k, v) in enumerate(metrics, start=2):
        c1 = summary.cell(row=i, column=1, value=k)
        c2 = summary.cell(row=i, column=2, value=v)
        c1.font = body_font
        c2.font = body_font
        c1.alignment = left
        c2.alignment = left
        if high_quality and i % 2 == 0:
            c1.fill = alt_fill
            c2.fill = alt_fill
    summary.column_dimensions["A"].width = 24
    summary.column_dimensions["B"].width = 28
    if image_path and os.path.isfile(image_path):
        try:
            from openpyxl.drawing.image import Image as XLImage  # type: ignore

            summary.add_image(XLImage(image_path), "D2")
        except Exception:
            pass
    wb.save(xlsx_path)


def write_audit_docx(
    docx_path: str,
    job_id: str,
    profile: Dict[str, Any],
    image_path: Optional[str] = None,
    params: Optional[Dict[str, Any]] = None,
    *,
    utc_now_str,
) -> None:
    from docx import Document  # type: ignore
    from docx.shared import Inches, Pt, RGBColor  # type: ignore

    theme = office_theme_settings(params)
    layout = office_layout_settings(params)
    high_quality = office_is_high_quality(params)
    primary_rgb = hex_to_rgb(str(theme.get("primary_hex", "1F4E78")))
    font_name = office_font_name(params)
    doc = Document()
    section = doc.sections[0]
    section.left_margin = Inches(0.9)
    section.right_margin = Inches(0.9)
    section.top_margin = Inches(0.8)
    section.bottom_margin = Inches(0.8)

    title = doc.add_heading(str(theme.get("report_title")), level=1)
    title.runs[0].font.size = Pt(20)
    for r in title.runs:
        r.font.name = font_name
        r.font.color.rgb = RGBColor(primary_rgb[0], primary_rgb[1], primary_rgb[2])

    subtitle = doc.add_paragraph(
        office_text(
            "本报告用于作业与辩论场景，展示清洗结果、质量指标和可复用结论。",
            "This report summarizes cleaned output, quality metrics, and reusable conclusions.",
            params,
        )
    )
    subtitle.runs[0].font.name = font_name
    subtitle.runs[0].font.size = Pt(10.5)

    meta = doc.add_table(rows=4, cols=2)
    meta.style = "Light List Accent 1"
    meta.cell(0, 0).text = office_text("任务ID", "Job ID", params)
    meta.cell(0, 1).text = job_id
    meta.cell(1, 0).text = office_text("流程步骤", "Step", params)
    meta.cell(1, 1).text = "cleaning"
    meta.cell(2, 0).text = office_text("生成时间", "Generated At", params)
    meta.cell(2, 1).text = utc_now_str()
    meta.cell(3, 0).text = office_text("状态", "Status", params)
    meta.cell(3, 1).text = office_text("完成", "DONE", params)
    meta.add_row().cells[0].text = office_text("主题", "Theme", params)
    meta.rows[-1].cells[1].text = str(theme.get("display_name", theme.get("name")))
    meta.add_row().cells[0].text = office_text("质量模式", "Quality Mode", params)
    meta.rows[-1].cells[1].text = office_quality_mode(params)

    doc.add_paragraph("")
    doc.add_heading(office_text("核心指标", "Core Metrics", params), level=2)
    metrics = doc.add_table(rows=1, cols=2)
    metrics.style = "Light Grid Accent 1"
    metrics.rows[0].cells[0].text = office_text("指标", "Metric", params)
    metrics.rows[0].cells[1].text = office_text("数值", "Value", params)
    core_items = [
        (office_text("行数", "Rows", params), profile.get("rows")),
        (office_text("列数", "Columns", params), profile.get("cols")),
        (office_text("金额总计", "Sum Amount", params), profile.get("sum_amount")),
        (office_text("金额最小值", "Min Amount", params), profile.get("min_amount")),
        (office_text("金额最大值", "Max Amount", params), profile.get("max_amount")),
        (office_text("金额均值", "Avg Amount", params), profile.get("avg_amount")),
    ]
    max_core = max(4, int(layout.get("docx_max_table_rows", 20)))
    for k, v in core_items[:max_core]:
        cells = metrics.add_row().cells
        cells[0].text = str(k)
        cells[1].text = str(v)

    quality = profile.get("quality") if isinstance(profile.get("quality"), dict) else {}
    doc.add_heading(office_text("质量摘要", "Quality Summary", params), level=2)
    q_items = [
        (office_text("输入行数", "Input Rows", params), quality.get("input_rows")),
        (office_text("输出行数", "Output Rows", params), quality.get("output_rows")),
        (office_text("无效行数", "Invalid Rows", params), quality.get("invalid_rows")),
        (office_text("过滤行数", "Filtered Rows", params), quality.get("filtered_rows")),
        (office_text("去重移除行数", "Duplicates Removed", params), quality.get("duplicate_rows_removed")),
    ]
    q = doc.add_table(rows=1, cols=2)
    q.style = "Light Grid Accent 1"
    q.rows[0].cells[0].text = office_text("质量指标", "Quality Metric", params)
    q.rows[0].cells[1].text = office_text("数值", "Value", params)
    for k, v in q_items:
        cells = q.add_row().cells
        cells[0].text = str(k)
        cells[1].text = str(v)

    doc.add_heading(office_text("结论要点", "Key Takeaways", params), level=2)
    invalid_rows = int(quality.get("invalid_rows", 0) or 0)
    filtered_rows = int(quality.get("filtered_rows", 0) or 0)
    rows_count = int(profile.get("rows", 0) or 0)
    takeaways = [
        office_text(
            f"最终可用数据 {rows_count} 行；可直接用于作业正文与辩论证据表。",
            f"Final usable rows: {rows_count}; ready for assignment writing and debate evidence.",
            params,
        ),
        office_text(
            f"无效行 {invalid_rows}，过滤行 {filtered_rows}；建议在附录中说明处理规则。",
            f"Invalid rows: {invalid_rows}, filtered rows: {filtered_rows}; explain data rules in appendix.",
            params,
        ),
        office_text(
            "图表与表格已按统一主题输出，可直接复制到课程汇报材料。",
            "Charts and tables are exported with a unified theme and can be reused directly.",
            params,
        ),
    ]
    for t in takeaways:
        p = doc.add_paragraph(style="List Bullet")
        run = p.add_run(t)
        run.font.name = font_name

    numeric_stats = profile.get("numeric_stats") if isinstance(profile.get("numeric_stats"), dict) else {}
    if numeric_stats:
        doc.add_heading(office_text("字段统计", "Field Statistics", params), level=2)
        stats_table = doc.add_table(rows=1, cols=5)
        stats_table.style = "Light Grid Accent 1"
        headers = [
            office_text("字段", "Field", params),
            office_text("总和", "Sum", params),
            office_text("最小值", "Min", params),
            office_text("最大值", "Max", params),
            office_text("均值", "Avg", params),
        ]
        for i, h in enumerate(headers):
            stats_table.rows[0].cells[i].text = h
        for field in sorted(numeric_stats.keys())[:12]:
            st = numeric_stats.get(field) if isinstance(numeric_stats.get(field), dict) else {}
            cells = stats_table.add_row().cells
            cells[0].text = str(field)
            cells[1].text = str(st.get("sum"))
            cells[2].text = str(st.get("min"))
            cells[3].text = str(st.get("max"))
            cells[4].text = str(st.get("avg"))
    if image_path and os.path.isfile(image_path):
        doc.add_heading(office_text("可视化说明", "Visual Explanation", params), level=2)
        doc.add_picture(image_path, width=Inches(6.2 if high_quality else 6.6))
    docx_apply_font(doc, font_name)
    doc.save(docx_path)


def write_deck_pptx(
    pptx_path: str,
    job_id: str,
    profile: Dict[str, Any],
    image_path: Optional[str] = None,
    params: Optional[Dict[str, Any]] = None,
    *,
    utc_now_str,
) -> None:
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore
        from pptx.dml.color import RGBColor  # type: ignore

        theme = office_theme_settings(params)
        layout = office_layout_settings(params)
        high_quality = office_is_high_quality(params)
        font_name = office_font_name(params)
        primary_rgb = hex_to_rgb(str(theme.get("primary_hex", "1F4E78")))
        secondary_rgb = hex_to_rgb(str(theme.get("secondary_hex", "2C6A9D")))
        accent_rgb = hex_to_rgb(str(theme.get("accent_hex", "F39C34")))
        bg_rgb = tuple(theme.get("bg_rgb") or (248, 250, 252))

        def add_bg(slide: Any) -> None:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(int(bg_rgb[0]), int(bg_rgb[1]), int(bg_rgb[2]))

        def add_top_bar(slide: Any) -> None:
            bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.24))
            bar.fill.solid()
            bar.fill.fore_color.rgb = RGBColor(primary_rgb[0], primary_rgb[1], primary_rgb[2])
            bar.line.fill.background()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        add_bg(slide)
        add_top_bar(slide)
        slide.shapes.title.text = str(theme.get("cover_title"))
        title_tf = slide.shapes.title.text_frame
        title_tf.paragraphs[0].font.size = Pt(38 if high_quality else 34)
        title_tf.paragraphs[0].font.color.rgb = RGBColor(primary_rgb[0], primary_rgb[1], primary_rgb[2])
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = (
                f"{office_text('任务', 'Job', params)} {job_id}\n"
                f"{office_text('主题', 'Theme', params)}: {theme.get('display_name', theme.get('name'))}\n"
                f"{office_text('生成时间', 'Generated at', params)}: {utc_now_str()}"
            )
            for p in slide.placeholders[1].text_frame.paragraphs:
                p.font.size = Pt(16 if high_quality else 14)
                p.font.color.rgb = RGBColor(secondary_rgb[0], secondary_rgb[1], secondary_rgb[2])

        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        add_bg(slide2)
        add_top_bar(slide2)
        slide2.shapes.title.text = office_text("关键指标", "Key Metrics", params)
        tb = slide2.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(8.7), Inches(4.9))
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        kpis = [
            f"{office_text('行数', 'Rows', params)}: {profile.get('rows')}",
            f"{office_text('列数', 'Columns', params)}: {profile.get('cols')}",
            f"{office_text('金额总计', 'Sum Amount', params)}: {profile.get('sum_amount')}",
            f"{office_text('金额均值', 'Avg Amount', params)}: {profile.get('avg_amount')}",
        ]
        for i, line in enumerate(kpis):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(27 if high_quality and i < 2 else 22 if i < 2 else 18)
            p.font.color.rgb = RGBColor(primary_rgb[0], primary_rgb[1], primary_rgb[2])

        quality = profile.get("quality") if isinstance(profile.get("quality"), dict) else {}
        slide3 = prs.slides.add_slide(prs.slide_layouts[1])
        add_bg(slide3)
        add_top_bar(slide3)
        slide3.shapes.title.text = office_text("数据质量", "Data Quality", params)
        qf = slide3.placeholders[1].text_frame
        qf.clear()
        qf.word_wrap = True
        q_items = [
            f"{office_text('输入行数', 'Input rows', params)}: {quality.get('input_rows', 0)}",
            f"{office_text('输出行数', 'Output rows', params)}: {quality.get('output_rows', 0)}",
            f"{office_text('无效行数', 'Invalid rows', params)}: {quality.get('invalid_rows', 0)}",
            f"{office_text('过滤行数', 'Filtered rows', params)}: {quality.get('filtered_rows', 0)}",
            f"{office_text('去重移除', 'Duplicates removed', params)}: {quality.get('duplicate_rows_removed', 0)}",
        ]
        max_items = max(4, int(layout.get("pptx_max_items", 6)))
        for i, line in enumerate(q_items[:max_items]):
            p = qf.paragraphs[0] if i == 0 else qf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(18 if high_quality else 16)
            p.font.color.rgb = RGBColor(secondary_rgb[0], secondary_rgb[1], secondary_rgb[2])

        numeric_stats = profile.get("numeric_stats") if isinstance(profile.get("numeric_stats"), dict) else {}
        if numeric_stats:
            slide4 = prs.slides.add_slide(prs.slide_layouts[1])
            add_bg(slide4)
            add_top_bar(slide4)
            slide4.shapes.title.text = office_text("字段统计(Top)", "Field Statistics (Top)", params)
            sf = slide4.placeholders[1].text_frame
            sf.clear()
            sf.word_wrap = True
            for i, field in enumerate(sorted(numeric_stats.keys())[:6]):
                st = numeric_stats.get(field) if isinstance(numeric_stats.get(field), dict) else {}
                txt = (
                    f"{field}: "
                    f"{office_text('总和', 'sum', params)}={st.get('sum')} "
                    f"{office_text('最小', 'min', params)}={st.get('min')} "
                    f"{office_text('最大', 'max', params)}={st.get('max')} "
                    f"{office_text('均值', 'avg', params)}={st.get('avg')}"
                )
                p = sf.paragraphs[0] if i == 0 else sf.add_paragraph()
                p.text = txt
                p.level = 0
                p.font.size = Pt(14 if high_quality else 13)
                p.font.color.rgb = RGBColor(secondary_rgb[0], secondary_rgb[1], secondary_rgb[2])

        slide_img = prs.slides.add_slide(prs.slide_layouts[5])
        add_bg(slide_img)
        add_top_bar(slide_img)
        slide_img.shapes.title.text = office_text("图表混排说明", "Table + Visual Mix", params)
        margin = Inches(0.4)
        top = Inches(1.2)
        gap = Inches(0.2)
        usable_w = prs.slide_width - (2 * margin)
        image_w = int(usable_w * 0.56)
        table_w = usable_w - image_w - gap
        content_h = prs.slide_height - top - Inches(0.4)
        table_data = [
            [office_text("指标", "Metric", params), office_text("数值", "Value", params)],
            [office_text("行数", "Rows", params), str(profile.get("rows"))],
            [office_text("列数", "Columns", params), str(profile.get("cols"))],
            [office_text("金额总计", "Sum", params), str(profile.get("sum_amount"))],
            [office_text("金额均值", "Avg", params), str(profile.get("avg_amount"))],
        ]
        if image_path and os.path.isfile(image_path):
            add_picture_fit(
                slide_img,
                image_path,
                margin,
                top,
                image_w,
                content_h,
            )
        table_h = min(Inches(4.4), content_h)
        table_x = margin + image_w + gap
        table = slide_img.shapes.add_table(len(table_data), 2, table_x, top, table_w, table_h).table
        for r, row_data in enumerate(table_data):
            for c, cell_text in enumerate(row_data):
                table.cell(r, c).text = cell_text

        for sl in prs.slides:
            try:
                from aiwf.office_style import ppt_apply_font
                ppt_apply_font(sl, font_name)
            except Exception:
                pass

        prs.save(pptx_path)
    except Exception as e:
        raise RuntimeError("python-pptx unavailable; cannot generate deck.pptx") from e
