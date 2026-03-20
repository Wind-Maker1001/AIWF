from __future__ import annotations

import os
from typing import Any, Dict, Optional

from aiwf.office_style import (
    add_picture_fit,
    hex_to_rgb,
    office_font_name,
    office_is_high_quality,
    office_layout_settings,
    office_text,
    office_theme_settings,
)


def write_deck_pptx(
    pptx_path: str,
    job_id: str,
    profile: Dict[str, Any],
    image_path: Optional[str] = None,
    params: Optional[Dict[str, Any]] = None,
    *,
    utc_now_str,
) -> None:
    try:
        from pptx import Presentation  # type: ignore
        from pptx.dml.color import RGBColor  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore

        theme = office_theme_settings(params)
        layout = office_layout_settings(params)
        high_quality = office_is_high_quality(params)
        font_name = office_font_name(params)
        primary_rgb = hex_to_rgb(str(theme.get("primary_hex", "1F4E78")))
        secondary_rgb = hex_to_rgb(str(theme.get("secondary_hex", "2C6A9D")))
        bg_rgb = tuple(theme.get("bg_rgb") or (248, 250, 252))

        def add_bg(slide: Any) -> None:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(int(bg_rgb[0]), int(bg_rgb[1]), int(bg_rgb[2]))

        def add_top_bar(slide: Any, prs: Any) -> None:
            bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.24))
            bar.fill.solid()
            bar.fill.fore_color.rgb = RGBColor(primary_rgb[0], primary_rgb[1], primary_rgb[2])
            bar.line.fill.background()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        add_bg(slide)
        add_top_bar(slide, prs)
        slide.shapes.title.text = str(theme.get("cover_title"))
        title_tf = slide.shapes.title.text_frame
        title_tf.paragraphs[0].font.size = Pt(38 if high_quality else 34)
        title_tf.paragraphs[0].font.color.rgb = RGBColor(primary_rgb[0], primary_rgb[1], primary_rgb[2])
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = (
                f"{office_text('任务', 'Job', params)} {job_id}\n"
                f"{office_text('主题', 'Theme', params)}: {theme.get('display_name', theme.get('name'))}\n"
                f"{office_text('生成时间', 'Generated at', params)}: {utc_now_str()}"
            )
            for paragraph in slide.placeholders[1].text_frame.paragraphs:
                paragraph.font.size = Pt(16 if high_quality else 14)
                paragraph.font.color.rgb = RGBColor(secondary_rgb[0], secondary_rgb[1], secondary_rgb[2])

        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        add_bg(slide2)
        add_top_bar(slide2, prs)
        slide2.shapes.title.text = office_text("关键指标", "Key Metrics", params)
        textbox = slide2.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(8.7), Inches(4.9))
        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        kpis = [
            f"{office_text('行数', 'Rows', params)}: {profile.get('rows')}",
            f"{office_text('列数', 'Columns', params)}: {profile.get('cols')}",
            f"{office_text('金额总计', 'Sum Amount', params)}: {profile.get('sum_amount')}",
            f"{office_text('金额均值', 'Avg Amount', params)}: {profile.get('avg_amount')}",
        ]
        for i, line in enumerate(kpis):
            paragraph = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            paragraph.text = line
            paragraph.font.size = Pt(27 if high_quality and i < 2 else 22 if i < 2 else 18)
            paragraph.font.color.rgb = RGBColor(primary_rgb[0], primary_rgb[1], primary_rgb[2])

        quality = profile.get("quality") if isinstance(profile.get("quality"), dict) else {}
        slide3 = prs.slides.add_slide(prs.slide_layouts[1])
        add_bg(slide3)
        add_top_bar(slide3, prs)
        slide3.shapes.title.text = office_text("数据质量", "Data Quality", params)
        quality_frame = slide3.placeholders[1].text_frame
        quality_frame.clear()
        quality_frame.word_wrap = True
        q_items = [
            f"{office_text('输入行数', 'Input rows', params)}: {quality.get('input_rows', 0)}",
            f"{office_text('输出行数', 'Output rows', params)}: {quality.get('output_rows', 0)}",
            f"{office_text('无效行数', 'Invalid rows', params)}: {quality.get('invalid_rows', 0)}",
            f"{office_text('过滤行数', 'Filtered rows', params)}: {quality.get('filtered_rows', 0)}",
            f"{office_text('去重移除', 'Duplicates removed', params)}: {quality.get('duplicate_rows_removed', 0)}",
        ]
        max_items = max(4, int(layout.get("pptx_max_items", 6)))
        for i, line in enumerate(q_items[:max_items]):
            paragraph = quality_frame.paragraphs[0] if i == 0 else quality_frame.add_paragraph()
            paragraph.text = line
            paragraph.level = 0
            paragraph.font.size = Pt(18 if high_quality else 16)
            paragraph.font.color.rgb = RGBColor(secondary_rgb[0], secondary_rgb[1], secondary_rgb[2])

        numeric_stats = profile.get("numeric_stats") if isinstance(profile.get("numeric_stats"), dict) else {}
        if numeric_stats:
            slide4 = prs.slides.add_slide(prs.slide_layouts[1])
            add_bg(slide4)
            add_top_bar(slide4, prs)
            slide4.shapes.title.text = office_text("字段统计(Top)", "Field Statistics (Top)", params)
            stats_frame = slide4.placeholders[1].text_frame
            stats_frame.clear()
            stats_frame.word_wrap = True
            for i, field in enumerate(sorted(numeric_stats.keys())[:6]):
                stats = numeric_stats.get(field) if isinstance(numeric_stats.get(field), dict) else {}
                text = (
                    f"{field}: "
                    f"{office_text('总和', 'sum', params)}={stats.get('sum')} "
                    f"{office_text('最小', 'min', params)}={stats.get('min')} "
                    f"{office_text('最大', 'max', params)}={stats.get('max')} "
                    f"{office_text('均值', 'avg', params)}={stats.get('avg')}"
                )
                paragraph = stats_frame.paragraphs[0] if i == 0 else stats_frame.add_paragraph()
                paragraph.text = text
                paragraph.level = 0
                paragraph.font.size = Pt(14 if high_quality else 13)
                paragraph.font.color.rgb = RGBColor(secondary_rgb[0], secondary_rgb[1], secondary_rgb[2])

        slide_img = prs.slides.add_slide(prs.slide_layouts[5])
        add_bg(slide_img)
        add_top_bar(slide_img, prs)
        slide_img.shapes.title.text = office_text("图表混排说明", "Table + Visual Mix", params)
        margin = Inches(0.4)
        top = Inches(1.2)
        gap = Inches(0.2)
        usable_w = prs.slide_width - (2 * margin)
        image_w = int(usable_w * 0.56)
        table_w = usable_w - image_w - gap
        content_h = prs.slide_height - top - Inches(0.4)
        table_data = [
            [office_text("指标", "Metric", params), office_text("数值", "Value", params)],
            [office_text("行数", "Rows", params), str(profile.get("rows"))],
            [office_text("列数", "Columns", params), str(profile.get("cols"))],
            [office_text("金额总计", "Sum", params), str(profile.get("sum_amount"))],
            [office_text("金额均值", "Avg", params), str(profile.get("avg_amount"))],
        ]
        if image_path and os.path.isfile(image_path):
            add_picture_fit(slide_img, image_path, margin, top, image_w, content_h)
        table_h = min(Inches(4.4), content_h)
        table_x = margin + image_w + gap
        table = slide_img.shapes.add_table(len(table_data), 2, table_x, top, table_w, table_h).table
        for row_idx, row_data in enumerate(table_data):
            for col_idx, cell_text in enumerate(row_data):
                table.cell(row_idx, col_idx).text = cell_text

        for slide_item in prs.slides:
            try:
                from aiwf.office_style import ppt_apply_font

                ppt_apply_font(slide_item, font_name)
            except Exception:
                pass

        prs.save(pptx_path)
    except Exception as exc:
        raise RuntimeError("python-pptx unavailable; cannot generate deck.pptx") from exc
