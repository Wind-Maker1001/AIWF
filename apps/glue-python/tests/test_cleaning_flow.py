import os
import json
import tempfile
import unittest
from unittest.mock import patch, Mock

os.environ.setdefault("NUMEXPR_MAX_THREADS", "8")
os.environ.setdefault("AIWF_ALLOW_EXTERNAL_JOB_ROOT", "1")

from aiwf.flows import cleaning
from aiwf.flows import cleaning_flow_materialization
from aiwf.flows.cleaning_precheck import run_cleaning_precheck
from aiwf.flows.cleaning_artifacts import (
    list_cleaning_artifact_details,
    list_cleaning_artifact_domains,
    materialize_accel_cleaning_artifacts,
    register_cleaning_artifact,
    unregister_cleaning_artifact,
)
from aiwf.flows.cleaning_advanced_quality import evaluate_advanced_quality
from aiwf.flows.office_artifacts import (
    list_office_artifact_details,
    list_office_artifact_domains,
    register_office_artifact,
    unregister_office_artifact,
)
from aiwf.governance_manual_reviews import list_manual_reviews
from aiwf.governance_quality_rule_sets import save_quality_rule_set


def make_job_context(job_root: str) -> dict[str, str]:
    job_root = os.path.normpath(job_root)
    return {
        "job_root": job_root,
        "stage_dir": os.path.join(job_root, "stage"),
        "artifacts_dir": os.path.join(job_root, "artifacts"),
        "evidence_dir": os.path.join(job_root, "evidence"),
    }


def with_job_context(job_root: str, **params):
    out = dict(params)
    out["job_context"] = make_job_context(job_root)
    return out


class CleaningFlowTests(unittest.TestCase):
    def test_evaluate_advanced_quality_supports_report_only_and_block_modes(self):
        rows = [
            {"amount": 1},
            {"amount": 2},
            {"amount": 3},
            {"amount": 4},
            {"amount": 1000},
        ]

        with patch("aiwf.flows.cleaning_advanced_quality.quality_check_v4_operator", return_value={"ok": False, "error": "offline"}):
            report_only = evaluate_advanced_quality(
                rows=rows,
                params_effective={
                    "quality_rules": {
                        "advanced_rules": {
                            "outlier_zscore": {"field": "amount", "max_z": 1.0},
                        }
                    }
                },
            )
            blocked = evaluate_advanced_quality(
                rows=rows,
                params_effective={
                    "quality_rules": {
                        "advanced_rules": {
                            "outlier_zscore": {"field": "amount", "max_z": 1.0},
                            "block_on_advanced_rules": True,
                        }
                    }
                },
            )

        self.assertTrue(report_only["enabled"])
        self.assertTrue(report_only["report_only"])
        self.assertFalse(report_only["blocked"])
        self.assertFalse(report_only["passed"])
        self.assertTrue(report_only["fallback_used"])
        self.assertFalse(blocked["report_only"])
        self.assertTrue(blocked["blocked"])

    def test_evaluate_advanced_quality_preserves_multiple_outlier_rules(self):
        rows = [
            {"amount": 1, "balance": 1},
            {"amount": 2, "balance": 2},
            {"amount": 3, "balance": 3},
            {"amount": 4, "balance": 4},
            {"amount": 1000, "balance": 1000},
        ]

        with patch("aiwf.flows.cleaning_advanced_quality.quality_check_v4_operator", return_value={"ok": False, "error": "offline"}):
            out = evaluate_advanced_quality(
                rows=rows,
                params_effective={
                    "quality_rules": {
                        "advanced_rules": {
                            "outlier_zscore": [
                                {"field": "amount", "max_z": 1.0},
                                {"field": "balance", "max_z": 1.0},
                            ]
                        }
                    }
                },
            )

        self.assertIsInstance(out["rules"]["outlier_zscore"], list)
        self.assertEqual(len(out["rules"]["outlier_zscore"]), 2)
        details = out["report"]["violations"][0]["details"]
        self.assertEqual({item["field"] for item in details}, {"amount", "balance"})

    def test_clean_rows_supports_dirty_row_filters_and_table_field_ops(self):
        out = cleaning._clean_rows(
            [
                {"amount": "¥ 12.5", "phone": "138 0013 8000", "account_no": " 6222-0011 8899 "},
                {"c1": "ID", "c2": "Amount", "c3": "Txn Date"},
                {"note": "note: imported from OCR"},
                {"amount": "Subtotal", "note": "subtotal line"},
                {"amount": "", "note": "   "},
            ],
            {
                "rules": {
                    "use_rust_v2": False,
                    "platform_mode": "generic",
                    "filters": [
                        {"op": "header_repeat_row", "header_values": ["id", "amount", "txn_date"], "min_matches": 2},
                        {"op": "note_row", "keywords": ["note"]},
                        {"op": "subtotal_row", "keywords": ["subtotal"]},
                        {"op": "blank_row"},
                    ],
                    "field_ops": [
                        {"field": "amount", "op": "strip_currency_symbol"},
                        {"field": "amount", "op": "parse_number"},
                        {"field": "amount", "op": "scale_by_header_unit", "unit": "万元"},
                        {"field": "phone", "op": "normalize_phone_cn"},
                        {"field": "account_no", "op": "normalize_account_no"},
                    ],
                }
            },
        )

        self.assertEqual(len(out["rows"]), 1)
        self.assertEqual(out["rows"][0]["phone"], "+8613800138000")
        self.assertEqual(out["rows"][0]["account_no"], "622200118899")
        self.assertEqual(out["rows"][0]["amount"], 125000.0)
        self.assertEqual(out["quality"]["filtered_rows"], 4)

    def test_office_theme_settings(self):
        t = cleaning._office_theme_settings({"office_theme": "debate"})
        self.assertEqual(t["name"], "debate")
        self.assertTrue("report_title" in t)
        t_assignment = cleaning._office_theme_settings({"office_theme": "assignment"})
        self.assertEqual(t_assignment["name"], "assignment")
        self.assertTrue("report_title" in t_assignment)
        t_en = cleaning._office_theme_settings({"office_theme": "debate", "office_lang": "en"})
        self.assertIn("Debate", t_en["report_title"])
        self.assertEqual(cleaning._office_quality_mode({}), "high")
        self.assertEqual(cleaning._office_quality_mode({"office_quality_mode": "standard"}), "standard")

    def test_build_profile_tracks_dynamic_columns_and_numeric_stats(self):
        profile = cleaning._build_profile(
            [
                {"id": 1, "amount": "10.25", "score": "2", "name": "alice"},
                {"id": 2, "amount": "1.75", "score": "3", "name": "bob"},
            ],
            {"input_rows": 2, "output_rows": 2, "invalid_rows": 0, "filtered_rows": 0, "duplicate_rows_removed": 0},
            "unit.test",
        )

        self.assertEqual(profile["cols"], 4)
        self.assertEqual(profile["sum_amount"], 12.0)
        self.assertEqual(profile["avg_amount"], 6.0)
        self.assertEqual(profile["numeric_stats"]["score"]["sum"], 5.0)
        self.assertEqual(profile["numeric_stats"]["amount"]["max"], 10.25)

    def test_build_profile_ignores_missing_amount_values(self):
        profile = cleaning._build_profile(
            [{"amount": 10}, {}],
            {"input_rows": 2, "output_rows": 2, "invalid_rows": 0, "filtered_rows": 0, "duplicate_rows_removed": 0},
            "unit.test",
        )

        self.assertEqual(profile["sum_amount"], 10.0)
        self.assertEqual(profile["min_amount"], 10.0)
        self.assertEqual(profile["max_amount"], 10.0)
        self.assertEqual(profile["avg_amount"], 10.0)

    def test_office_writers_produce_rich_outputs(self):
        with tempfile.TemporaryDirectory() as tmp:
            xlsx_path = os.path.join(tmp, "fin.xlsx")
            docx_path = os.path.join(tmp, "audit.docx")
            pptx_path = os.path.join(tmp, "deck.pptx")
            img_path = os.path.join(tmp, "summary_visual.png")
            rows = [{"id": 1, "amount": 10.2, "name": "alice"}, {"id": 2, "amount": 20.3, "name": "bob"}]
            profile = cleaning._build_profile(
                rows,
                {"input_rows": 2, "output_rows": 2, "invalid_rows": 0, "filtered_rows": 0, "duplicate_rows_removed": 0},
                "unit.test",
            )
            cleaning._write_profile_illustration_png(img_path, profile)
            cleaning._write_fin_xlsx(xlsx_path, rows, img_path)
            cleaning._write_audit_docx(docx_path, "job-x", profile, img_path)
            cleaning._write_deck_pptx(pptx_path, "job-x", profile, img_path)

            from openpyxl import load_workbook  # type: ignore
            from docx import Document  # type: ignore

            wb = load_workbook(xlsx_path, read_only=True, data_only=True)
            self.assertEqual(wb.sheetnames, ["detail", "summary"])
            wb.close()
            doc = Document(docx_path)
            self.assertTrue("报告" in doc.paragraphs[0].text or "Report" in doc.paragraphs[0].text)
            self.assertTrue(os.path.getsize(pptx_path) > 0)

            from pptx import Presentation  # type: ignore

            prs = Presentation(pptx_path)
            self.assertGreaterEqual(len(prs.slides), 5)
            sw = prs.slide_width
            sh = prs.slide_height
            for s in prs.slides:
                for shape in s.shapes:
                    self.assertLessEqual(shape.left + shape.width, sw)
                    self.assertLessEqual(shape.top + shape.height, sh)

    def test_office_writers_support_english(self):
        with tempfile.TemporaryDirectory() as tmp:
            docx_path = os.path.join(tmp, "audit_en.docx")
            rows = [{"id": 1, "amount": 10.2}]
            profile = cleaning._build_profile(
                rows,
                {"input_rows": 1, "output_rows": 1, "invalid_rows": 0, "filtered_rows": 0, "duplicate_rows_removed": 0},
                "unit.test",
            )
            cleaning._write_audit_docx(docx_path, "job-en", profile, params={"office_lang": "en", "office_theme": "academic"})
            from docx import Document  # type: ignore

            doc = Document(docx_path)
            self.assertEqual(doc.paragraphs[0].text, "Academic Data Processing Report")

    def test_run_cleaning_includes_registered_custom_office_artifact(self):
        with tempfile.TemporaryDirectory() as tmp:
            local_job_root = os.path.join(tmp, "job")

            def write_valid_parquet(path, rows):
                with open(path, "wb") as f:
                    f.write(b"PAR1dataPAR1")

            def write_bin(path, *args, **kwargs):
                with open(path, "wb") as f:
                    f.write(b"BIN")
                return True

            def write_text_summary(context, output_path):
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(f"rows={context.office_profile.get('rows')}")

            register_office_artifact(
                "text_summary",
                artifact_id="text_summary_001",
                kind="txt",
                filename="summary.txt",
                path_key="summary_txt_path",
                sha_key="sha_summary_txt",
                writer=write_text_summary,
                domain="custom-office",
                domain_metadata={"label": "Custom Office", "backend": "extension", "builtin": False},
            )
            try:
                with patch("aiwf.flows.cleaning._base_step_start"), patch(
                    "aiwf.flows.cleaning._base_artifact_upsert"
                ), patch("aiwf.flows.cleaning._base_step_done"), patch(
                    "aiwf.flows.cleaning._base_step_fail"
                ), patch(
                    "aiwf.flows.cleaning._try_accel_cleaning",
                    return_value={"attempted": True, "ok": False, "error": "accel unavailable"},
                ), patch(
                    "aiwf.flows.cleaning._write_cleaned_parquet", side_effect=write_valid_parquet
                ), patch(
                    "aiwf.flows.cleaning._write_profile_illustration_png", side_effect=write_bin
                ), patch(
                    "aiwf.flows.cleaning._write_fin_xlsx", side_effect=write_bin
                ), patch(
                    "aiwf.flows.cleaning._write_audit_docx", side_effect=write_bin
                ), patch(
                    "aiwf.flows.cleaning._write_deck_pptx", side_effect=write_bin
                ):
                    out = cleaning.run_cleaning(
                        job_id="job-custom-office",
                        actor="test",
                        params=with_job_context(local_job_root, rows=[{"id": 1, "amount": 10.0}]),
                    )
                    details = {item["name"]: item for item in list_office_artifact_details()}
                    domains = list_office_artifact_domains()
            finally:
                unregister_office_artifact("text_summary")

            txt_artifacts = [a for a in out["artifacts"] if a["kind"] == "txt"]
            self.assertEqual(len(txt_artifacts), 1)
            self.assertTrue(txt_artifacts[0]["path"].endswith("summary.txt"))
            self.assertEqual(details["text_summary"]["domain"], "custom-office")
            self.assertTrue(any(item["name"] == "custom-office" for item in domains))

    def test_run_cleaning_includes_registered_custom_core_artifact(self):
        with tempfile.TemporaryDirectory() as tmp:
            local_job_root = os.path.join(tmp, "job")

            def write_valid_parquet(path, rows):
                with open(path, "wb") as f:
                    f.write(b"PAR1dataPAR1")

            def custom_meta_path(context):
                return os.path.join(context.evidence_dir, "meta.txt")

            def write_custom_meta(context, output_path):
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(f"rows={len(context.rows)}")

            register_cleaning_artifact(
                "meta_text",
                artifact_id="meta_text_001",
                kind="txt",
                path_key="meta_txt_path",
                sha_key="sha_meta_txt",
                local_path_resolver=custom_meta_path,
                local_writer=write_custom_meta,
                domain="custom-core",
                domain_metadata={"label": "Custom Core", "backend": "extension", "builtin": False},
            )
            try:
                with patch("aiwf.flows.cleaning._base_step_start"), patch(
                    "aiwf.flows.cleaning._base_artifact_upsert"
                ), patch("aiwf.flows.cleaning._base_step_done"), patch(
                    "aiwf.flows.cleaning._base_step_fail"
                ), patch(
                    "aiwf.flows.cleaning._try_accel_cleaning",
                    return_value={"attempted": True, "ok": False, "error": "accel unavailable"},
                ), patch(
                    "aiwf.flows.cleaning._write_cleaned_parquet", side_effect=write_valid_parquet
                ):
                    out = cleaning.run_cleaning(
                        job_id="job-custom-core",
                        actor="test",
                        params=with_job_context(local_job_root, rows=[{"id": 1, "amount": 10.0}]),
                    )
                    details = {item["name"]: item for item in list_cleaning_artifact_details()}
                    domains = list_cleaning_artifact_domains()
            finally:
                unregister_cleaning_artifact("meta_text")

            txt_artifacts = [a for a in out["artifacts"] if a["artifact_id"] == "meta_text_001"]
            self.assertEqual(len(txt_artifacts), 1)
            self.assertTrue(txt_artifacts[0]["path"].endswith("meta.txt"))
            self.assertEqual(details["meta_text"]["domain"], "custom-core")
            self.assertTrue(any(item["name"] == "custom-core" for item in domains))

    def test_run_cleaning_can_disable_office_outputs(self):
        with tempfile.TemporaryDirectory() as tmp:
            local_job_root = os.path.join(tmp, "job")

            def write_valid_parquet(path, rows):
                with open(path, "wb") as f:
                    f.write(b"PAR1dataPAR1")

            with patch("aiwf.flows.cleaning._base_step_start"), patch(
                "aiwf.flows.cleaning._base_artifact_upsert"
            ), patch("aiwf.flows.cleaning._base_step_done"), patch(
                "aiwf.flows.cleaning._base_step_fail"
            ), patch(
                "aiwf.flows.cleaning._try_accel_cleaning",
                return_value={"attempted": True, "ok": False, "error": "accel unavailable"},
            ), patch(
                "aiwf.flows.cleaning._write_cleaned_parquet", side_effect=write_valid_parquet
            ), patch(
                "aiwf.flows.cleaning._write_profile_illustration_png", side_effect=AssertionError("office disabled")
            ), patch(
                "aiwf.flows.cleaning._write_fin_xlsx", side_effect=AssertionError("office disabled")
            ), patch(
                "aiwf.flows.cleaning._write_audit_docx", side_effect=AssertionError("office disabled")
            ), patch(
                "aiwf.flows.cleaning._write_deck_pptx", side_effect=AssertionError("office disabled")
            ):
                out = cleaning.run_cleaning(
                    job_id="job-no-office",
                    actor="test",
                    params=with_job_context(
                        local_job_root,
                        office_outputs_enabled=False,
                        rows=[{"id": 1, "amount": 10.0}],
                    ),
                )

            artifact_kinds = {a["kind"] for a in out["artifacts"]}
            self.assertFalse({"xlsx", "docx", "pptx"} & artifact_kinds)

    def test_run_cleaning_can_disable_optional_core_artifacts(self):
        with tempfile.TemporaryDirectory() as tmp:
            local_job_root = os.path.join(tmp, "job")

            def write_valid_parquet(path, rows):
                with open(path, "wb") as f:
                    f.write(b"PAR1dataPAR1")

            with patch("aiwf.flows.cleaning._base_step_start"), patch(
                "aiwf.flows.cleaning._base_artifact_upsert"
            ), patch("aiwf.flows.cleaning._base_step_done"), patch(
                "aiwf.flows.cleaning._base_step_fail"
            ), patch(
                "aiwf.flows.cleaning._try_accel_cleaning",
                return_value={"attempted": True, "ok": False, "error": "accel unavailable"},
            ), patch(
                "aiwf.flows.cleaning._write_cleaned_parquet", side_effect=write_valid_parquet
            ):
                out = cleaning.run_cleaning(
                    job_id="job-core-select",
                    actor="test",
                    params=with_job_context(
                        local_job_root,
                        office_outputs_enabled=False,
                        disabled_core_artifacts=["csv", "json", "jsonl"],
                        rows=[{"id": 1, "amount": 10.0}],
                    ),
                )

            artifact_kinds = [a["kind"] for a in out["artifacts"]]
            self.assertEqual(artifact_kinds, ["parquet"])

    def test_run_cleaning_rejects_disabling_required_parquet_artifact(self):
        with tempfile.TemporaryDirectory() as tmp:
            local_job_root = os.path.join(tmp, "job")

            with patch("aiwf.flows.cleaning._base_step_start"), patch(
                "aiwf.flows.cleaning._base_artifact_upsert"
            ), patch("aiwf.flows.cleaning._base_step_done"), patch(
                "aiwf.flows.cleaning._base_step_fail"
            ), patch(
                "aiwf.flows.cleaning._try_accel_cleaning",
                return_value={"attempted": True, "ok": False, "error": "accel unavailable"},
            ):
                with self.assertRaises(RuntimeError) as ctx:
                    cleaning.run_cleaning(
                        job_id="job-disable-required",
                        actor="test",
                        params=with_job_context(
                            local_job_root,
                            office_outputs_enabled=False,
                            disabled_core_artifacts=["parquet"],
                            rows=[{"id": 1, "amount": 10.0}],
                        ),
                    )
            self.assertIn("required cleaning artifact disabled", str(ctx.exception))

    def test_run_cleaning_supports_nested_artifact_selection(self):
        with tempfile.TemporaryDirectory() as tmp:
            local_job_root = os.path.join(tmp, "job")

            def write_valid_parquet(path, rows):
                with open(path, "wb") as f:
                    f.write(b"PAR1dataPAR1")

            with patch("aiwf.flows.cleaning._base_step_start"), patch(
                "aiwf.flows.cleaning._base_artifact_upsert"
            ), patch("aiwf.flows.cleaning._base_step_done"), patch(
                "aiwf.flows.cleaning._base_step_fail"
            ), patch(
                "aiwf.flows.cleaning._try_accel_cleaning",
                return_value={"attempted": True, "ok": False, "error": "accel unavailable"},
            ), patch(
                "aiwf.flows.cleaning._write_cleaned_parquet", side_effect=write_valid_parquet
            ), patch(
                "aiwf.flows.cleaning._write_profile_illustration_png", side_effect=AssertionError("office disabled")
            ), patch(
                "aiwf.flows.cleaning._write_fin_xlsx", side_effect=AssertionError("office disabled")
            ), patch(
                "aiwf.flows.cleaning._write_audit_docx", side_effect=AssertionError("office disabled")
            ), patch(
                "aiwf.flows.cleaning._write_deck_pptx", side_effect=AssertionError("office disabled")
            ):
                out = cleaning.run_cleaning(
                    job_id="job-nested-artifacts",
                    actor="test",
                    params=with_job_context(
                        local_job_root,
                        artifact_selection={
                            "office": {"enabled": False},
                            "core": {"enabled": False},
                        },
                        rows=[{"id": 1, "amount": 10.0}],
                    ),
                )

            artifact_kinds = [a["kind"] for a in out["artifacts"]]
            self.assertEqual(artifact_kinds, ["parquet"])

    def test_run_cleaning_applies_quality_rule_set_and_writes_quality_summary_artifacts(self):
        with tempfile.TemporaryDirectory() as tmp:
            local_job_root = os.path.join(tmp, "job")

            def write_valid_parquet(path, rows):
                with open(path, "wb") as f:
                    f.write(b"PAR1dataPAR1")

            with patch.dict("os.environ", {"AIWF_GOVERNANCE_ROOT": tmp}, clear=False):
                save_quality_rule_set(
                    {
                        "id": "finance_default",
                        "name": "Finance Default",
                        "rules": {"min_output_rows": 1},
                    }
                )
                with patch("aiwf.flows.cleaning._base_step_start"), patch(
                    "aiwf.flows.cleaning._base_artifact_upsert"
                ), patch("aiwf.flows.cleaning._base_step_done"), patch(
                    "aiwf.flows.cleaning._base_step_fail"
                ), patch(
                    "aiwf.flows.cleaning._try_accel_cleaning",
                    return_value={"attempted": True, "ok": False, "error": "accel unavailable"},
                ), patch(
                    "aiwf.flows.cleaning._write_cleaned_parquet", side_effect=write_valid_parquet
                ):
                    out = cleaning.run_cleaning(
                        job_id="job-quality-summary",
                        actor="test",
                        params=with_job_context(
                            local_job_root,
                            office_outputs_enabled=False,
                            quality_rule_set_id="finance_default",
                            rules={"use_rust_v2": False},
                            rows=[{"id": "1", "amount": "10"}, {"id": "bad", "amount": "20"}],
                        ),
                    )

            self.assertEqual(out["quality_summary"]["rule_set_provenance"]["resolved_id"], "finance_default")
            self.assertEqual(out["quality_summary"]["engine_path"]["execution_mode"], "python_legacy")
            quality_summary_artifact = next(a for a in out["artifacts"] if a["artifact_id"] == "quality_summary_json_001")
            rejections_artifact = next(a for a in out["artifacts"] if a["artifact_id"] == "rejections_jsonl_001")
            with open(quality_summary_artifact["path"], "r", encoding="utf-8") as f:
                quality_summary_payload = json.load(f)
            self.assertEqual(quality_summary_payload["schema_version"], "cleaning_quality_summary.v1")
            self.assertEqual(quality_summary_payload["rule_set_provenance"]["resolved_id"], "finance_default")
            with open(rejections_artifact["path"], "r", encoding="utf-8") as f:
                rejection_lines = [json.loads(line) for line in f if line.strip()]
            self.assertTrue(any(item["reason_category"] == "required_missing" for item in rejection_lines))
            self.assertTrue(all("reason_code" in item for item in rejection_lines))

    def test_run_cleaning_auto_enqueues_manual_review_for_survivorship_tie(self):
        with tempfile.TemporaryDirectory() as tmp:
            local_job_root = os.path.join(tmp, "job")

            def write_valid_parquet(path, rows):
                with open(path, "wb") as f:
                    f.write(b"PAR1dataPAR1")

            with patch.dict("os.environ", {"AIWF_GOVERNANCE_ROOT": tmp}, clear=False), patch(
                "aiwf.flows.cleaning._write_cleaned_parquet", side_effect=write_valid_parquet
            ):
                out = cleaning.run_cleaning(
                    job_id="job-manual-review",
                    actor="test",
                    params=with_job_context(
                        local_job_root,
                        local_standalone=True,
                        office_outputs_enabled=False,
                        cleaning_template="customer_contact_v1",
                        rules={
                            "use_rust_v2": False,
                            "platform_mode": "generic",
                            "deduplicate_by": ["phone"],
                            "survivorship": {
                                "keys": ["phone"],
                                "tie_breaker": "last",
                            },
                        },
                        rows=[
                            {"phone": "13800138000", "customer_name": "Alice"},
                            {"phone": "13800138000", "customer_name": "Alicia"},
                        ],
                    ),
                )
                queued = list_manual_reviews()

            self.assertTrue(out["quality_summary"]["review_analysis"]["review_required"])
            self.assertTrue(out["quality_summary"]["manual_review_queue"]["auto_enqueued"])
            self.assertEqual(len(queued), 1)
            self.assertEqual(queued[0]["review_key"], "cleaning::duplicate_key_risk::1")

    def test_run_cleaning_enqueues_manual_review_when_risky_duplicate_is_outside_sample(self):
        with tempfile.TemporaryDirectory() as tmp:
            local_job_root = os.path.join(tmp, "job")

            def write_valid_parquet(path, rows):
                with open(path, "wb") as f:
                    f.write(b"PAR1dataPAR1")

            with patch.dict("os.environ", {"AIWF_GOVERNANCE_ROOT": tmp}, clear=False), patch(
                "aiwf.flows.cleaning._write_cleaned_parquet", side_effect=write_valid_parquet
            ):
                out = cleaning.run_cleaning(
                    job_id="job-manual-review-sampled",
                    actor="test",
                    params=with_job_context(
                        local_job_root,
                        local_standalone=True,
                        office_outputs_enabled=False,
                        audit_sample_limit=1,
                        rules={
                            "use_rust_v2": False,
                            "platform_mode": "generic",
                            "deduplicate_by": ["phone"],
                            "survivorship": {
                                "keys": ["phone"],
                                "prefer_non_null_fields": ["customer_name"],
                                "prefer_latest_fields": ["biz_date"],
                                "tie_breaker": "last",
                            },
                        },
                        rows=[
                            {"phone": "safe", "customer_name": "", "biz_date": "2026-01-01"},
                            {"phone": "safe", "customer_name": "Alice", "biz_date": "2026-01-02"},
                            {"phone": "risk", "customer_name": "Bob", "biz_date": "2026-01-01"},
                            {"phone": "risk", "customer_name": "Bobby", "biz_date": "2026-01-01"},
                        ],
                    ),
                )
                queued = list_manual_reviews()

            duplicate_risk = out["quality_summary"]["review_analysis"]["duplicate_key_risk"]
            self.assertEqual(duplicate_risk["duplicate_rows_removed"], 2)
            self.assertTrue(duplicate_risk["review_required"])
            self.assertTrue(out["quality_summary"]["manual_review_queue"]["auto_enqueued"])
            self.assertEqual(len(queued), 1)
            self.assertEqual(queued[0]["review_key"], "cleaning::duplicate_key_risk::1")

    def test_clean_rows_scales_english_header_units_in_python_path(self):
        out = cleaning._clean_rows(
            [{"amount": "12.5"}],
            {
                "rules": {
                    "use_rust_v2": False,
                    "platform_mode": "generic",
                    "field_ops": [
                        {"field": "amount", "op": "parse_number"},
                        {"field": "amount", "op": "scale_by_header_unit", "unit": "million"},
                    ],
                }
            },
        )

        self.assertEqual(out["rows"][0]["amount"], 12500000.0)

    def test_clean_rows_parse_number_supports_full_width_comma_in_python_path(self):
        out = cleaning._clean_rows(
            [{"amount": "1，234.56"}],
            {
                "rules": {
                    "use_rust_v2": False,
                    "platform_mode": "generic",
                    "field_ops": [
                        {"field": "amount", "op": "parse_number"},
                    ],
                }
            },
        )

        self.assertEqual(out["rows"][0]["amount"], 1234.56)

    def test_clean_rows_parse_number_supports_accounting_notation_in_python_path(self):
        out = cleaning._clean_rows(
            [{"amount": "(1,234.50)"}, {"amount": "1,234.50-"}],
            {
                "rules": {
                    "use_rust_v2": False,
                    "platform_mode": "generic",
                    "field_ops": [
                        {"field": "amount", "op": "parse_number"},
                    ],
                }
            },
        )

        self.assertEqual(out["rows"][0]["amount"], -1234.5)
        self.assertEqual(out["rows"][1]["amount"], -1234.5)

    def test_clean_rows_sign_amount_from_direction_field_for_bank_rows(self):
        out = cleaning._clean_rows(
            [
                {"amount": "120.50", "txn_type": "借"},
                {"amount": "300.00", "txn_type": "贷"},
                {"amount": "999.00", "debit_amount": "120.50", "credit_amount": "", "txn_type": "贷"},
            ],
            {
                "rules": {
                    "use_rust_v2": False,
                    "platform_mode": "generic",
                    "field_ops": [
                        {"field": "amount", "op": "parse_number"},
                        {"field": "debit_amount", "op": "parse_number"},
                        {"field": "credit_amount", "op": "parse_number"},
                        {"field": "amount", "op": "sign_amount_from_debit_credit", "direction_field": "txn_type"},
                    ],
                }
            },
        )

        self.assertEqual(out["rows"][0]["amount"], -120.5)
        self.assertEqual(out["rows"][1]["amount"], 300.0)
        self.assertEqual(out["rows"][2]["amount"], -120.5)

    def test_evaluate_advanced_quality_preserves_custom_direction_field_for_bank_semantics(self):
        out = evaluate_advanced_quality(
            rows=[
                {
                    "account_no": "6222-0001",
                    "txn_date": "2026-03-01",
                    "drcr": "debit",
                    "amount": "120.50",
                    "balance": "10000.00",
                }
            ],
            params_effective={
                "canonical_profile": "bank_statement",
                "quality_rules": {
                    "advanced_rules": {
                        "bank_statement_semantics": {
                            "direction_field": "drcr",
                            "block_on_semantic_conflicts": True,
                        }
                    }
                },
            },
        )

        self.assertTrue(out["enabled"])
        self.assertFalse(out["passed"])
        self.assertEqual(
            out["semantic_checks"]["summary"]["counts"]["signed_amount_conflict"],
            1,
        )

    def test_evaluate_advanced_quality_handles_non_numeric_row_index_for_bank_semantics(self):
        out = evaluate_advanced_quality(
            rows=[
                {
                    "row_index": "A1",
                    "account_no": "6222-0001",
                    "txn_date": "2026-03-01",
                    "debit_amount": "120.50",
                    "credit_amount": "",
                    "amount": "999.00",
                    "balance": "10000.00",
                }
            ],
            params_effective={
                "canonical_profile": "bank_statement",
                "quality_rules": {
                    "advanced_rules": {
                        "bank_statement_semantics": {
                            "block_on_semantic_conflicts": True,
                        }
                    }
                },
            },
        )

        self.assertTrue(out["enabled"])
        self.assertEqual(out["semantic_checks"]["items"][0]["row_index"], 1)

    def test_evaluate_advanced_quality_reports_bank_statement_semantic_conflicts(self):
        out = evaluate_advanced_quality(
            rows=[
                {
                    "account_no": "6222-0001",
                    "txn_date": "2026-03-01",
                    "debit_amount": "120.50",
                    "credit_amount": "",
                    "amount": "999.00",
                    "balance": "10000.00",
                }
            ],
            params_effective={
                "canonical_profile": "bank_statement",
                "quality_rules": {
                    "advanced_rules": {
                        "bank_statement_semantics": {
                            "block_on_semantic_conflicts": True,
                        }
                    }
                },
            },
        )

        self.assertTrue(out["enabled"])
        self.assertTrue(out["blocked"])
        self.assertFalse(out["passed"])
        self.assertEqual(
            out["semantic_checks"]["summary"]["counts"]["signed_amount_conflict"],
            1,
        )

    def test_run_cleaning_precheck_warns_on_bank_semantic_conflict(self):
        payload = run_cleaning_precheck(
            params={"canonical_profile": "bank_statement"},
            extract_payload={
                "rows": [
                    {
                        "account_no": "6222-0001",
                        "txn_date": "2026-03-01",
                        "debit_amount": "120.50",
                        "credit_amount": "",
                        "amount": "999.00",
                        "balance": "10000.00",
                    }
                ],
                "header_mapping": [],
                "candidate_profiles": [
                    {
                        "profile": "bank_statement",
                        "recommended": True,
                        "score": 0.95,
                        "required_coverage": 1.0,
                        "recommended_template_id": "bank_statement_v1",
                    }
                ],
                "quality_decisions": [],
                "sample_rows": [],
                "quality_blocked": False,
                "blocked_reason_codes": [],
            },
        )

        self.assertEqual(payload["precheck_action"], "warn")
        self.assertTrue(payload["review_required"])
        self.assertTrue(any(item["kind"] == "signed_amount_conflict" for item in payload["review_items"]))

    def test_run_cleaning_precheck_blocks_on_bank_semantic_conflict_when_enabled(self):
        payload = run_cleaning_precheck(
            params={
                "canonical_profile": "bank_statement",
                "quality_rules": {
                    "advanced_rules": {
                        "bank_statement_semantics": {
                            "block_on_semantic_conflicts": True,
                        }
                    }
                },
            },
            extract_payload={
                "rows": [
                    {
                        "account_no": "6222-0001",
                        "txn_date": "2026-03-01",
                        "debit_amount": "120.50",
                        "credit_amount": "",
                        "amount": "999.00",
                        "balance": "10000.00",
                    }
                ],
                "header_mapping": [],
                "candidate_profiles": [
                    {
                        "profile": "bank_statement",
                        "recommended": True,
                        "score": 0.95,
                        "required_coverage": 1.0,
                        "recommended_template_id": "bank_statement_v1",
                    }
                ],
                "quality_decisions": [],
                "sample_rows": [],
                "quality_blocked": False,
                "blocked_reason_codes": [],
            },
        )

        self.assertEqual(payload["precheck_action"], "block")
        self.assertIn("signed_amount_conflict", payload["blocking_reason_codes"])

    def test_run_cleaning_precheck_warns_on_bank_balance_gap(self):
        payload = run_cleaning_precheck(
            params={"canonical_profile": "bank_statement"},
            extract_payload={
                "rows": [
                    {
                        "account_no": "6222-0001",
                        "txn_date": "2026-03-01",
                        "amount": "-100.00",
                        "balance": "900.00",
                        "ref_no": "TXN-001",
                    },
                    {
                        "account_no": "6222-0001",
                        "txn_date": "2026-03-02",
                        "amount": "200.00",
                        "balance": "950.00",
                        "ref_no": "TXN-002",
                    },
                ],
                "header_mapping": [],
                "candidate_profiles": [
                    {
                        "profile": "bank_statement",
                        "recommended": True,
                        "score": 0.95,
                        "required_coverage": 1.0,
                        "recommended_template_id": "bank_statement_v1",
                    }
                ],
                "quality_decisions": [],
                "sample_rows": [],
                "quality_blocked": False,
                "blocked_reason_codes": [],
            },
        )

        self.assertEqual(payload["precheck_action"], "warn")
        self.assertTrue(payload["review_required"])
        self.assertTrue(any(item["kind"] == "balance_gap" for item in payload["review_items"]))

    def test_run_cleaning_precheck_blocks_on_bank_balance_gap_when_enabled(self):
        payload = run_cleaning_precheck(
            params={
                "canonical_profile": "bank_statement",
                "quality_rules": {
                    "advanced_rules": {
                        "bank_statement_semantics": {
                            "block_on_semantic_conflicts": True,
                        }
                    }
                },
            },
            extract_payload={
                "rows": [
                    {
                        "account_no": "6222-0001",
                        "txn_date": "2026-03-01",
                        "amount": "-100.00",
                        "balance": "900.00",
                        "ref_no": "TXN-001",
                    },
                    {
                        "account_no": "6222-0001",
                        "txn_date": "2026-03-02",
                        "amount": "200.00",
                        "balance": "950.00",
                        "ref_no": "TXN-002",
                    },
                ],
                "header_mapping": [],
                "candidate_profiles": [
                    {
                        "profile": "bank_statement",
                        "recommended": True,
                        "score": 0.95,
                        "required_coverage": 1.0,
                        "recommended_template_id": "bank_statement_v1",
                    }
                ],
                "quality_decisions": [],
                "sample_rows": [],
                "quality_blocked": False,
                "blocked_reason_codes": [],
            },
        )

        self.assertEqual(payload["precheck_action"], "block")
        self.assertIn("balance_gap", payload["blocking_reason_codes"])

    def test_clean_rows_bank_semantics_uses_cleaned_rows_after_filtering(self):
        out = cleaning._clean_rows(
            [
                {
                    "row_index": 1,
                    "account_no": "6222-0001",
                    "txn_date": "2026/03/01",
                    "amount": "-100.00",
                    "balance": "900.00",
                    "ref_no": "TXN-001",
                },
                {
                    "row_index": 2,
                    "account_no": "6222-0001",
                    "txn_date": "2026/03/01",
                    "amount": "50.00",
                    "balance": "950.00",
                    "remark": "subtotal row",
                    "ref_no": "TXN-SUBTOTAL",
                },
                {
                    "row_index": 3,
                    "account_no": "6222-0001",
                    "txn_date": "2026/03/02",
                    "amount": "200.00",
                    "balance": "1100.00",
                    "ref_no": "TXN-002",
                },
            ],
            {
                "canonical_profile": "bank_statement",
                "rules": {
                    "use_rust_v2": False,
                    "platform_mode": "generic",
                    "filters": [{"op": "subtotal_row"}],
                    "field_ops": [
                        {"field": "account_no", "op": "normalize_account_no"},
                        {"field": "txn_date", "op": "parse_date"},
                        {"field": "amount", "op": "parse_number"},
                        {"field": "balance", "op": "parse_number"},
                    ],
                },
                "quality_rules": {
                    "advanced_rules": {
                        "bank_statement_semantics": {
                            "block_on_semantic_conflicts": True,
                        }
                    }
                },
            },
        )

        self.assertEqual(out["quality"]["filtered_rows"], 1)
        self.assertEqual(
            out["execution_audit"]["semantic_checks"]["summary"]["counts"]["balance_gap"],
            0,
        )
        self.assertFalse(
            any(item["kind"] == "balance_gap" for item in out["review_analysis"]["review_items"])
        )

    def test_run_cleaning_precheck_uses_cleaned_rows_for_bank_semantics_after_filtering(self):
        payload = run_cleaning_precheck(
            params={
                "canonical_profile": "bank_statement",
                "rules": {
                    "use_rust_v2": False,
                    "platform_mode": "generic",
                    "filters": [{"op": "subtotal_row"}],
                    "field_ops": [
                        {"field": "account_no", "op": "normalize_account_no"},
                        {"field": "txn_date", "op": "parse_date"},
                        {"field": "amount", "op": "parse_number"},
                        {"field": "balance", "op": "parse_number"},
                    ],
                },
                "quality_rules": {
                    "advanced_rules": {
                        "bank_statement_semantics": {
                            "block_on_semantic_conflicts": True,
                        }
                    }
                },
            },
            extract_payload={
                "rows": [
                    {
                        "row_index": 1,
                        "account_no": "6222-0001",
                        "txn_date": "2026/03/01",
                        "amount": "-100.00",
                        "balance": "900.00",
                        "ref_no": "TXN-001",
                    },
                    {
                        "row_index": 2,
                        "account_no": "6222-0001",
                        "txn_date": "2026/03/01",
                        "amount": "50.00",
                        "balance": "950.00",
                        "remark": "subtotal row",
                        "ref_no": "TXN-SUBTOTAL",
                    },
                    {
                        "row_index": 3,
                        "account_no": "6222-0001",
                        "txn_date": "2026/03/02",
                        "amount": "200.00",
                        "balance": "1100.00",
                        "ref_no": "TXN-002",
                    },
                ],
                "header_mapping": [],
                "candidate_profiles": [
                    {
                        "profile": "bank_statement",
                        "recommended": True,
                        "score": 0.95,
                        "required_coverage": 1.0,
                        "recommended_template_id": "bank_statement_v1",
                    }
                ],
                "quality_decisions": [],
                "sample_rows": [],
                "quality_blocked": False,
                "blocked_reason_codes": [],
            },
        )

        self.assertEqual(payload["precheck_action"], "allow")
        self.assertFalse(any(item["kind"] == "balance_gap" for item in payload["review_items"]))

    def test_run_cleaning_auto_enqueues_manual_review_for_bank_balance_gap(self):
        with tempfile.TemporaryDirectory() as tmp:
            local_job_root = os.path.join(tmp, "job")

            def write_valid_parquet(path, rows):
                with open(path, "wb") as f:
                    f.write(b"PAR1dataPAR1")

            with patch.dict("os.environ", {"AIWF_GOVERNANCE_ROOT": tmp}, clear=False), patch(
                "aiwf.flows.cleaning._write_cleaned_parquet", side_effect=write_valid_parquet
            ):
                out = cleaning.run_cleaning(
                    job_id="job-bank-balance-review",
                    actor="test",
                    params=with_job_context(
                        local_job_root,
                        local_standalone=True,
                        office_outputs_enabled=False,
                        cleaning_template="bank_statement_v1",
                        rules={"use_rust_v2": False},
                        rows=[
                            {
                                "account_no": "6222-0001",
                                "txn_date": "2026-03-01",
                                "debit_amount": "100.00",
                                "credit_amount": "",
                                "amount": "-100.00",
                                "balance": "900.00",
                                "counterparty_name": "张三",
                                "ref_no": "TXN-001",
                            },
                            {
                                "account_no": "6222-0001",
                                "txn_date": "2026-03-02",
                                "debit_amount": "",
                                "credit_amount": "200.00",
                                "amount": "200.00",
                                "balance": "950.00",
                                "counterparty_name": "李四",
                                "ref_no": "TXN-002",
                            },
                        ],
                    ),
                )
                queued = list_manual_reviews()

            self.assertTrue(out["quality_summary"]["review_analysis"]["review_required"])
            self.assertTrue(any(item["kind"] == "balance_gap" for item in out["quality_summary"]["review_analysis"]["review_items"]))
            self.assertTrue(out["quality_summary"]["manual_review_queue"]["auto_enqueued"])
            self.assertEqual(len(queued), 1)

    def test_materialize_accel_outputs_blocks_when_bank_balance_gap_blocks(self):
        raw_rows = [
            {
                "account_no": "6222-0001",
                "txn_date": "2026-03-01",
                "amount": "-100.00",
                "balance": "900.00",
                "ref_no": "TXN-001",
            },
            {
                "account_no": "6222-0001",
                "txn_date": "2026-03-02",
                "amount": "200.00",
                "balance": "950.00",
                "ref_no": "TXN-002",
            },
        ]
        with self.assertRaises(cleaning.CleaningGuardrailError) as ctx:
            cleaning_flow_materialization.materialize_accel_outputs(
                params_effective={
                    "canonical_profile": "bank_statement",
                    "quality_rules": {
                        "advanced_rules": {
                            "bank_statement_semantics": {
                                "block_on_semantic_conflicts": True,
                            }
                        }
                    },
                },
                accel_outputs={"cleaned_parquet": {"path": "D:/tmp/fake.parquet"}},
                accel_profile={"quality": {"output_rows": 2}, "quality_gate": {}},
                sha256_file=lambda _path: "sha",
                local_rows=list(raw_rows),
                local_profile={"quality": {"output_rows": 2}, "quality_gate": {}},
                local_execution={},
                preprocess_result={},
                input_rows=list(raw_rows),
            )
        self.assertEqual(ctx.exception.error_code, "advanced_quality_blocked")
        self.assertEqual(
            ctx.exception.details["advanced_quality"]["semantic_checks"]["summary"]["counts"]["balance_gap"],
            1,
        )

    def test_materialize_accel_outputs_uses_cleaned_rows_for_bank_semantics(self):
        with tempfile.TemporaryDirectory() as tmp:
            accel_stage = os.path.join(tmp, "accel")
            os.makedirs(accel_stage, exist_ok=True)
            cleaned_parquet = os.path.join(accel_stage, "cleaned.parquet")
            with open(cleaned_parquet, "wb") as f:
                f.write(b"PAR1dataPAR1")
            job_root = os.path.join(tmp, "job")
            raw_rows = [
                {
                    "account_no": "6222-0001",
                    "txn_date": "2026-03-01",
                    "amount": "-100.00",
                    "balance": "900.00",
                    "ref_no": "TXN-001",
                },
                {
                    "account_no": "6222-0001",
                    "txn_date": "2026-03-01",
                    "amount": "50.00",
                    "balance": "950.00",
                    "remark": "subtotal row",
                    "ref_no": "TXN-SUBTOTAL",
                },
                {
                    "account_no": "6222-0001",
                    "txn_date": "2026-03-02",
                    "amount": "200.00",
                    "balance": "1100.00",
                    "ref_no": "TXN-002",
                },
            ]
            cleaned_rows = [
                {
                    "account_no": "62220001",
                    "txn_date": "2026-03-01",
                    "amount": -100.0,
                    "balance": 900.0,
                    "ref_no": "TXN-001",
                },
                {
                    "account_no": "62220001",
                    "txn_date": "2026-03-02",
                    "amount": 200.0,
                    "balance": 1100.0,
                    "ref_no": "TXN-002",
                },
            ]
            out = cleaning_flow_materialization.materialize_accel_outputs(
                params_effective={
                    "canonical_profile": "bank_statement",
                    "office_outputs_enabled": False,
                    "job_context": make_job_context(job_root),
                    "quality_rules": {
                        "advanced_rules": {
                            "bank_statement_semantics": {
                                "block_on_semantic_conflicts": True,
                            }
                        }
                    },
                },
                accel_outputs={"cleaned_parquet": {"path": cleaned_parquet, "sha256": "parquet-sha"}},
                accel_profile={"quality": {"output_rows": 2}, "quality_gate": {}},
                sha256_file=lambda _path: "sha",
                local_rows=list(cleaned_rows),
                local_profile={"quality": {"output_rows": 2}, "quality_gate": {}},
                local_execution={},
                preprocess_result={},
                input_rows=list(raw_rows),
            )

            self.assertEqual(
                out["quality_summary"]["semantic_checks"]["summary"]["counts"]["balance_gap"],
                0,
            )
            self.assertEqual(
                out["quality_summary"]["advanced_quality"]["semantic_checks"]["summary"]["counts"]["balance_gap"],
                0,
            )
            self.assertEqual(
                out["execution"]["semantic_checks"]["summary"]["counts"]["balance_gap"],
                0,
            )
            self.assertFalse(out["quality_summary"]["advanced_quality"]["blocked"])

    def test_materialize_local_outputs_uses_cleaned_rows_for_bank_semantics(self):
        with tempfile.TemporaryDirectory() as tmp:
            stage_dir = os.path.join(tmp, "stage")
            artifacts_dir = os.path.join(tmp, "artifacts")
            evidence_dir = os.path.join(tmp, "evidence")
            os.makedirs(stage_dir, exist_ok=True)
            os.makedirs(artifacts_dir, exist_ok=True)
            os.makedirs(evidence_dir, exist_ok=True)

            raw_rows = [
                {
                    "account_no": "6222-0001",
                    "txn_date": "2026-03-01",
                    "amount": "-100.00",
                    "balance": "900.00",
                    "ref_no": "TXN-001",
                },
                {
                    "account_no": "6222-0001",
                    "txn_date": "2026-03-01",
                    "amount": "50.00",
                    "balance": "950.00",
                    "remark": "subtotal row",
                    "ref_no": "TXN-SUBTOTAL",
                },
                {
                    "account_no": "6222-0001",
                    "txn_date": "2026-03-02",
                    "amount": "200.00",
                    "balance": "1100.00",
                    "ref_no": "TXN-002",
                },
            ]
            cleaned_rows = [
                {
                    "account_no": "62220001",
                    "txn_date": "2026-03-01",
                    "amount": -100.0,
                    "balance": 900.0,
                    "ref_no": "TXN-001",
                },
                {
                    "account_no": "62220001",
                    "txn_date": "2026-03-02",
                    "amount": 200.0,
                    "balance": 1100.0,
                    "ref_no": "TXN-002",
                },
            ]

            def write_csv(path, rows):
                with open(path, "w", encoding="utf-8", newline="\n") as f:
                    f.write("account_no,txn_date,amount,balance,ref_no\n")
                    for row in rows:
                        f.write(
                            f"{row.get('account_no','')},{row.get('txn_date','')},{row.get('amount','')},{row.get('balance','')},{row.get('ref_no','')}\n"
                        )

            def write_parquet(path, rows):
                with open(path, "wb") as f:
                    f.write(b"PAR1dataPAR1")

            def write_profile_json(path, profile, params_effective):
                with open(path, "w", encoding="utf-8") as f:
                    json.dump(profile, f, ensure_ascii=False)

            out = cleaning_flow_materialization.materialize_local_outputs(
                job_id="job-local-semantic-cleaned-rows",
                stage_dir=stage_dir,
                artifacts_dir=artifacts_dir,
                evidence_dir=evidence_dir,
                params_effective={
                    "canonical_profile": "bank_statement",
                    "quality_rules": {
                        "advanced_rules": {
                            "bank_statement_semantics": {
                                "block_on_semantic_conflicts": True,
                            }
                        }
                    },
                },
                input_rows=list(raw_rows),
                rows=list(cleaned_rows),
                quality={"input_rows": 3, "output_rows": 2, "filtered_rows": 1},
                execution_report={},
                source="python",
                preprocess_result={},
                apply_quality_gates=lambda quality, _params: {"passed": True, "quality": dict(quality)},
                to_bool=lambda value, default=False: bool(default if value is None else value),
                rule_param=lambda _params, _key, default=None: default,
                require_local_parquet_dependencies=lambda _params: None,
                write_cleaned_csv=write_csv,
                write_cleaned_parquet=write_parquet,
                is_valid_parquet_file=lambda _path: True,
                local_parquet_strict_enabled=lambda _params: True,
                build_profile=lambda rows, quality, source: {"rows": len(rows), "source": source, "quality": dict(quality)},
                write_profile_json=write_profile_json,
                sha256_file=lambda _path: "sha",
                materialize_office_outputs_fn=lambda **_kwargs: {},
            )

            self.assertEqual(
                out["quality_summary"]["semantic_checks"]["summary"]["counts"]["balance_gap"],
                0,
            )
            self.assertEqual(
                out["quality_summary"]["advanced_quality"]["semantic_checks"]["summary"]["counts"]["balance_gap"],
                0,
            )
            self.assertEqual(
                out["execution"]["semantic_checks"]["summary"]["counts"]["balance_gap"],
                0,
            )
            self.assertFalse(out["quality_summary"]["advanced_quality"]["blocked"])

    def test_materialize_accel_outputs_blocks_when_advanced_quality_blocks(self):
        with self.assertRaises(cleaning.CleaningGuardrailError) as ctx:
            cleaning_flow_materialization.materialize_accel_outputs(
                params_effective={
                    "quality_rules": {
                        "advanced_rules": {
                            "outlier_zscore": {"field": "amount", "max_z": 1.0},
                            "block_on_advanced_rules": True,
                        }
                    }
                },
                accel_outputs={"cleaned_parquet": {"path": "D:/tmp/fake.parquet"}},
                accel_profile={"quality": {"output_rows": 5}, "quality_gate": {}},
                sha256_file=lambda _path: "sha",
                local_rows=[
                    {"amount": 1},
                    {"amount": 2},
                    {"amount": 3},
                    {"amount": 4},
                    {"amount": 1000},
                ],
                local_profile={"quality": {"output_rows": 5}, "quality_gate": {}},
                local_execution={},
                preprocess_result={},
                input_rows=[],
            )
        self.assertEqual(ctx.exception.error_code, "advanced_quality_blocked")

    def test_run_cleaning_quality_rule_set_keeps_explicit_quality_rule_overrides(self):
        with tempfile.TemporaryDirectory() as tmp:
            local_job_root = os.path.join(tmp, "job")

            def write_valid_parquet(path, rows):
                with open(path, "wb") as f:
                    f.write(b"PAR1dataPAR1")

            with patch.dict("os.environ", {"AIWF_GOVERNANCE_ROOT": tmp}, clear=False):
                save_quality_rule_set(
                    {
                        "id": "finance_override",
                        "name": "Finance Override",
                        "rules": {"min_output_rows": 2},
                    }
                )
                with patch("aiwf.flows.cleaning._base_step_start"), patch(
                    "aiwf.flows.cleaning._base_artifact_upsert"
                ), patch("aiwf.flows.cleaning._base_step_done"), patch(
                    "aiwf.flows.cleaning._base_step_fail"
                ), patch(
                    "aiwf.flows.cleaning._try_accel_cleaning",
                    return_value={"attempted": True, "ok": False, "error": "accel unavailable"},
                ), patch(
                    "aiwf.flows.cleaning._write_cleaned_parquet", side_effect=write_valid_parquet
                ):
                    out = cleaning.run_cleaning(
                        job_id="job-quality-override",
                        actor="test",
                        params=with_job_context(
                            local_job_root,
                            office_outputs_enabled=False,
                            quality_rule_set_id="finance_override",
                            quality_rules={"min_output_rows": 1},
                            rules={"use_rust_v2": False},
                            rows=[{"id": "1", "amount": "10"}],
                        ),
                    )

            self.assertEqual(out["profile"]["quality_gate"]["min_output_rows"], 1)
            self.assertIn("min_output_rows", out["quality_summary"]["rule_set_provenance"]["override_keys"])

    def test_run_cleaning_accel_outputs_still_emit_quality_summary_and_rejections(self):
        with tempfile.TemporaryDirectory() as tmp:
            local_job_root = os.path.join(tmp, "job")
            accel_stage = os.path.join(tmp, "accel")
            os.makedirs(accel_stage, exist_ok=True)
            cleaned_parquet = os.path.join(accel_stage, "cleaned.parquet")
            with open(cleaned_parquet, "wb") as f:
                f.write(b"PAR1dataPAR1")

            with patch("aiwf.flows.cleaning._base_step_start"), patch(
                "aiwf.flows.cleaning._base_artifact_upsert"
            ), patch("aiwf.flows.cleaning._base_step_done"), patch(
                "aiwf.flows.cleaning._base_step_fail"
            ), patch(
                "aiwf.flows.cleaning._try_accel_cleaning",
                return_value={
                    "attempted": True,
                    "ok": True,
                    "response": {
                        "outputs": {
                            "cleaned_parquet": {"path": cleaned_parquet, "sha256": "parquet-sha"},
                        },
                        "profile": {"rows": 1, "cols": 2},
                    },
                },
            ):
                out = cleaning.run_cleaning(
                    job_id="job-accel-quality-summary",
                    actor="test",
                    params=with_job_context(
                        local_job_root,
                        office_outputs_enabled=False,
                        rules={"use_rust_v2": False},
                        rows=[{"id": "1", "amount": "10"}, {"id": "bad", "amount": "20"}],
                    ),
                )

            self.assertFalse(out["accel"]["used_fallback"])
            self.assertEqual(out["quality_summary"]["engine_path"]["execution_mode"], "accel_operator")
            self.assertEqual(out["quality_summary"]["engine_path"]["row_transform_engine"], "python")
            self.assertEqual(out["quality_summary"]["engine_path"]["materialization_engine"], "legacy_accel_cleaning")
            artifact_ids = {artifact["artifact_id"] for artifact in out["artifacts"]}
            self.assertIn("quality_summary_json_001", artifact_ids)
            self.assertIn("rejections_jsonl_001", artifact_ids)

    def test_quality_summary_engine_path_legacy_cleaning_flag_matches_accel_block(self):
        with tempfile.TemporaryDirectory() as tmp:
            local_job_root = os.path.join(tmp, "job")
            accel_stage = os.path.join(tmp, "accel")
            os.makedirs(accel_stage, exist_ok=True)
            cleaned_parquet = os.path.join(accel_stage, "cleaned.parquet")
            with open(cleaned_parquet, "wb") as f:
                f.write(b"PAR1dataPAR1")

            with patch("aiwf.flows.cleaning._base_step_start"), patch(
                "aiwf.flows.cleaning._base_artifact_upsert"
            ), patch("aiwf.flows.cleaning._base_step_done"), patch(
                "aiwf.flows.cleaning._base_step_fail"
            ), patch(
                "aiwf.flows.cleaning._try_accel_cleaning",
                return_value={
                    "attempted": True,
                    "ok": True,
                    "response": {
                        "outputs": {
                            "cleaned_parquet": {"path": cleaned_parquet, "sha256": "parquet-sha"},
                        },
                        "profile": {"rows": 1, "cols": 2},
                    },
                },
            ):
                out = cleaning.run_cleaning(
                    job_id="job-legacy-flag",
                    actor="test",
                    params=with_job_context(
                        local_job_root,
                        office_outputs_enabled=False,
                        rules={"use_rust_v2": False},
                        rows=[{"id": "1", "amount": "10"}],
                    ),
                )

            self.assertEqual(
                out["quality_summary"]["engine_path"]["legacy_cleaning_operator_used"],
                out["accel"]["legacy_cleaning_used"],
            )

    def test_validate_cleaning_rules_ok(self):
        res = cleaning.validate_cleaning_rules(
            {
                "rules": {
                    "platform_mode": "generic",
                    "casts": {"amount": "float"},
                    "filters": [{"field": "amount", "op": "gte", "value": 0}],
                    "max_invalid_ratio": 0.2,
                    "max_required_missing_ratio": 0.1,
                    "deduplicate_keep": "last",
                }
            }
        )
        self.assertTrue(res["ok"])
        self.assertEqual(res["errors"], [])

    def test_validate_cleaning_rules_error(self):
        res = cleaning.validate_cleaning_rules(
            {
                "rules": {
                    "platform_mode": "abc",
                    "amount_round_digits": 99,
                    "max_invalid_ratio": 2,
                    "max_required_missing_ratio": 3,
                    "filters": [{"op": "bad"}],
                }
            }
        )
        self.assertFalse(res["ok"])
        self.assertGreaterEqual(len(res["errors"]), 3)

    def test_validate_cleaning_rules_warn_unknown(self):
        res = cleaning.validate_cleaning_rules({"rules": {"unknown_x": 1}})
        self.assertTrue(res["ok"])
        self.assertTrue(any("unknown rule keys" in w for w in res["warnings"]))

    def test_validate_cleaning_rules_supports_artifact_selection_object(self):
        res = cleaning.validate_cleaning_rules(
            {
                "rules": {"max_invalid_rows": 1},
                "artifact_selection": {
                    "office": {"enabled": False},
                    "core": {"exclude": ["csv"]},
                },
            }
        )
        self.assertTrue(res["ok"])

        bad = cleaning.validate_cleaning_rules(
            {
                "artifact_selection": {
                    "office": {"enabled": "no"},
                    "core": {"exclude": "csv"},
                }
            }
        )
        self.assertFalse(bad["ok"])

        bad_unknown = cleaning.validate_cleaning_rules(
            {
                "artifact_selection": {
                    "office": {"include": ["missing_office"]},
                    "core": {"exclude": ["missing_core"]},
                }
            }
        )
        self.assertFalse(bad_unknown["ok"])
        self.assertTrue(any("unknown artifacts" in err for err in bad_unknown["errors"]))

    def test_quality_gates_fail_on_invalid_rows(self):
        quality = {"input_rows": 10, "output_rows": 8, "invalid_rows": 2, "filtered_rows": 0}
        with self.assertRaises(RuntimeError) as ctx:
            cleaning._apply_quality_gates(quality, {"max_invalid_rows": 1})
        self.assertIn("max_invalid_rows", str(ctx.exception))

    def test_quality_gates_fail_on_invalid_ratio(self):
        quality = {"input_rows": 10, "output_rows": 8, "invalid_rows": 2, "filtered_rows": 0}
        with self.assertRaises(RuntimeError) as ctx:
            cleaning._apply_quality_gates(quality, {"rules": {"max_invalid_ratio": "0.1"}})
        self.assertIn("max_invalid_ratio", str(ctx.exception))

    def test_quality_gates_pass(self):
        quality = {"input_rows": 10, "output_rows": 8, "invalid_rows": 1, "filtered_rows": 1}
        out = cleaning._apply_quality_gates(
            quality,
            {
                "rules": {
                    "max_invalid_rows": 2,
                    "max_filtered_rows": 2,
                    "min_output_rows": 5,
                    "max_invalid_ratio": "0.2",
                }
            },
        )
        self.assertTrue(out["evaluated"])

    def test_clean_rows_generic_rules_pipeline(self):
        raw_rows = [
            {"name": " Alice ", "amt": "12.4", "city": "shanghai"},
            {"name": "bob", "amt": "9.9", "city": "beijing"},
            {"name": "alice", "amt": "100.1", "city": "shanghai"},
        ]
        cleaned = cleaning._clean_rows(
            raw_rows,
            {
                "rules": {
                    "platform_mode": "generic",
                    "rename_map": {"amt": "amount"},
                    "casts": {"amount": "float"},
                    "trim_strings": True,
                    "lowercase_fields": ["name"],
                    "filters": [{"field": "amount", "op": "gte", "value": 10}],
                    "deduplicate_by": ["name"],
                    "deduplicate_keep": "first",
                    "sort_by": [{"field": "amount", "order": "desc"}],
                }
            },
        )
        self.assertEqual(cleaned["quality"]["output_rows"], 1)
        self.assertEqual(cleaned["rows"][0]["name"], "alice")
        self.assertEqual(cleaned["rows"][0]["amount"], 12.4)

    def test_clean_rows_generic_rules_supports_bank_statement_computed_amount(self):
        raw_rows = [
            {
                "账号": "62220001",
                "交易日期": "2026/03/01",
                "借方金额": "120.50",
                "贷方金额": "0",
                "余额": "1000.00",
                "流水号": "TXN-1",
            }
        ]
        cleaned = cleaning._clean_rows(
            raw_rows,
            {
                "rules": {
                    "use_rust_v2": False,
                    "platform_mode": "generic",
                    "rename_map": {
                        "账号": "account_no",
                        "交易日期": "txn_date",
                        "借方金额": "debit_amount",
                        "贷方金额": "credit_amount",
                        "余额": "balance",
                        "流水号": "ref_no",
                    },
                    "casts": {
                        "debit_amount": "float",
                        "credit_amount": "float",
                        "balance": "float",
                    },
                    "required_fields": ["account_no", "txn_date"],
                    "computed_fields": {
                        "amount": "sub($credit_amount,$debit_amount)",
                    },
                },
            },
        )
        self.assertEqual(cleaned["rows"][0]["account_no"], "62220001")
        self.assertEqual(cleaned["rows"][0]["txn_date"], "2026/03/01")
        self.assertEqual(cleaned["rows"][0]["amount"], -120.5)

    def test_try_accel_cleaning_sends_params_payload(self):
        fake_resp = Mock()
        fake_resp.status_code = 200
        fake_resp.json.return_value = {"ok": True}
        with patch("requests.post", return_value=fake_resp) as post:
            cleaning._try_accel_cleaning(
                params={"rows": [{"id": 1, "amount": 2}], "rules": {"max_amount": 10}},
                job_id="j1",
                step_id="cleaning",
                actor="local",
                ruleset_version="v1",
                input_uri="in",
                output_uri="out",
            )
            kwargs = post.call_args.kwargs
            sent = kwargs["json"]
            self.assertIn("params", sent)
            self.assertIn("rows", sent["params"])
            self.assertIn("rules", sent["params"])
    def test_clean_rows_supports_declarative_rules_block(self):
        raw_rows = [
            {"ID": "1", "AMT": "10.5"},
            {"ID": "1", "AMT": "11.5"},
            {"ID": "2", "AMT": "-1"},
        ]
        cleaned = cleaning._clean_rows(
            raw_rows,
            {
                "rules": {
                    "id_field": "ID",
                    "amount_field": "AMT",
                    "drop_negative_amount": True,
                    "deduplicate_by_id": True,
                    "deduplicate_keep": "first",
                    "amount_round_digits": 0,
                }
            },
        )
        self.assertEqual(cleaned["rows"], [{"id": 1, "amount": 11.0}])
        self.assertEqual(cleaned["quality"]["rule_hits"]["filtered_negative"], 1)

    def test_clean_rows_applies_filters_dedup_and_rounding(self):
        raw_rows = [
            {"id": "1", "amount": "100.126"},
            {"id": "2", "amount": "-5"},
            {"id": "bad", "amount": "3"},
            {"id": "1", "amount": "$120.225"},
            {"id": "3", "amount": "999"},
        ]
        cleaned = cleaning._clean_rows(
            raw_rows,
            {
                "drop_negative_amount": True,
                "max_amount": 500,
                "deduplicate_by_id": True,
                "deduplicate_keep": "last",
                "amount_round_digits": 2,
            },
        )
        rows = cleaned["rows"]
        quality = cleaned["quality"]

        self.assertEqual(rows, [{"id": 1, "amount": 120.23}])
        self.assertEqual(quality["input_rows"], 5)
        self.assertEqual(quality["output_rows"], 1)
        self.assertEqual(quality["invalid_rows"], 1)
        self.assertEqual(quality["filtered_rows"], 2)
        self.assertEqual(quality["duplicate_rows_removed"], 1)

    def test_load_rows_from_input_csv_path(self):
        with tempfile.TemporaryDirectory() as tmp:
            csv_path = os.path.join(tmp, "in.csv")
            with open(csv_path, "w", encoding="utf-8", newline="\n") as f:
                f.write("id,amount\n10,88.1\n11,91.2\n")

            rows, source = cleaning._load_raw_rows({"input_csv_path": csv_path}, tmp)
            self.assertEqual(source, csv_path)
            self.assertEqual(len(rows), 2)
            self.assertEqual(rows[0]["id"], "10")
            self.assertEqual(rows[1]["amount"], "91.2")

    def test_load_rows_rejects_empty_explicit_rows(self):
        with self.assertRaisesRegex(RuntimeError, "params.rows is empty"):
            cleaning._load_raw_rows({"rows": []}, None)

    def test_load_rows_rejects_missing_input_csv_path(self):
        with tempfile.TemporaryDirectory() as tmp:
            missing_path = os.path.join(tmp, "missing.csv")
            with self.assertRaisesRegex(FileNotFoundError, "input csv file not found"):
                cleaning._load_raw_rows({"input_csv_path": missing_path}, tmp)

    def test_load_rows_rejects_header_only_input_csv_path(self):
        with tempfile.TemporaryDirectory() as tmp:
            csv_path = os.path.join(tmp, "empty.csv")
            with open(csv_path, "w", encoding="utf-8", newline="\n") as f:
                f.write("id,amount\n")
            with self.assertRaisesRegex(RuntimeError, "has no data rows"):
                cleaning._load_raw_rows({"input_csv_path": csv_path}, tmp)

    def test_clean_rows_large_input_quality_counts(self):
        raw_rows = [{"id": i, "amount": i} for i in range(1, 2001)]
        raw_rows.append({"id": "bad", "amount": "1"})
        raw_rows.append({"id": 2, "amount": -3})
        cleaned = cleaning._clean_rows(
            raw_rows,
            {
                "drop_negative_amount": True,
                "max_amount": 1500,
            },
        )
        quality = cleaned["quality"]
        self.assertEqual(quality["input_rows"], 2002)
        self.assertEqual(quality["invalid_rows"], 1)
        self.assertEqual(quality["filtered_rows"], 501)
        self.assertEqual(quality["output_rows"], 1500)

    def test_is_valid_parquet_file_checks_magic_bytes(self):
        with tempfile.TemporaryDirectory() as tmp:
            ok_path = os.path.join(tmp, "ok.parquet")
            bad_path = os.path.join(tmp, "bad.parquet")

            with open(ok_path, "wb") as f:
                f.write(b"PAR1abcdPAR1")
            with open(bad_path, "wb") as f:
                f.write(b"PARQUET_PLACEHOLDER\n")

            self.assertTrue(cleaning._is_valid_parquet_file(ok_path))
            self.assertFalse(cleaning._is_valid_parquet_file(bad_path))
            self.assertFalse(cleaning._is_valid_parquet_file(os.path.join(tmp, "missing.parquet")))

    def test_run_cleaning_falls_back_when_accel_parquet_invalid(self):
        with tempfile.TemporaryDirectory() as tmp:
            accel_stage = os.path.join(tmp, "accel")
            os.makedirs(accel_stage, exist_ok=True)
            accel_parquet = os.path.join(accel_stage, "cleaned.parquet")
            with open(accel_parquet, "wb") as f:
                f.write(b"PARQUET_PLACEHOLDER\n")

            local_job_root = os.path.join(tmp, "job")

            def write_local_csv(path, rows):
                with open(path, "w", encoding="utf-8") as f:
                    f.write("id,amount\n1,100\n2,200\n")
                return {"rows": 2, "cols": 2}

            def write_local_parquet(path, rows):
                with open(path, "wb") as f:
                    f.write(b"PAR1dataPAR1")

            def write_local_bin(path, *args, **kwargs):
                with open(path, "wb") as f:
                    f.write(b"BIN")

            def write_local_json(path, profile, params):
                with open(path, "w", encoding="utf-8") as f:
                    f.write("{}")

            with patch("aiwf.flows.cleaning._base_step_start"), patch(
                "aiwf.flows.cleaning._base_artifact_upsert"
            ), patch("aiwf.flows.cleaning._base_step_done"), patch(
                "aiwf.flows.cleaning._base_step_fail"
            ), patch(
                "aiwf.flows.cleaning._try_accel_cleaning",
                return_value={
                    "attempted": True,
                    "ok": True,
                    "response": {
                        "outputs": {
                            "cleaned_csv": {"path": accel_parquet, "sha256": ""},
                            "cleaned_parquet": {"path": accel_parquet, "sha256": ""},
                            "profile_json": {"path": accel_parquet, "sha256": ""},
                            "xlsx_fin": {"path": accel_parquet, "sha256": ""},
                            "audit_docx": {"path": accel_parquet, "sha256": ""},
                            "deck_pptx": {"path": accel_parquet, "sha256": ""},
                        },
                        "profile": {"rows": 2, "cols": 2},
                    },
                },
            ), patch(
                "aiwf.flows.cleaning._write_cleaned_csv", side_effect=write_local_csv
            ), patch(
                "aiwf.flows.cleaning._write_cleaned_parquet", side_effect=write_local_parquet
            ), patch(
                "aiwf.flows.cleaning._write_fin_xlsx", side_effect=write_local_bin
            ), patch(
                "aiwf.flows.cleaning._write_audit_docx", side_effect=write_local_bin
            ), patch(
                "aiwf.flows.cleaning._write_deck_pptx", side_effect=write_local_bin
            ), patch(
                "aiwf.flows.cleaning._write_profile_json", side_effect=write_local_json
            ):
                out = cleaning.run_cleaning(
                    job_id="job-1",
                    actor="test",
                    params=with_job_context(local_job_root, rows=[{"id": 1, "amount": 100.0}]),
                )

            self.assertTrue(out["ok"])
            self.assertTrue(out["accel"]["attempted"])
            self.assertTrue(out["accel"]["used_fallback"])
            self.assertIn("invalid parquet", str(out["accel"]["validation_error"]))
            parquet_artifact = [a for a in out["artifacts"] if a["kind"] == "parquet"][0]
            self.assertTrue(parquet_artifact["path"].startswith(local_job_root))
            self.assertNotEqual(parquet_artifact["path"], accel_parquet)

    def test_run_cleaning_uses_accel_office_outputs_when_accel_succeeds(self):
        with tempfile.TemporaryDirectory() as tmp:
            local_job_root = os.path.join(tmp, "job")
            accel_stage = os.path.join(tmp, "accel_stage")
            accel_artifacts = os.path.join(tmp, "accel_artifacts")
            accel_evidence = os.path.join(tmp, "accel_evidence")
            os.makedirs(accel_stage, exist_ok=True)
            os.makedirs(accel_artifacts, exist_ok=True)
            os.makedirs(accel_evidence, exist_ok=True)

            cleaned_csv = os.path.join(accel_stage, "cleaned.csv")
            cleaned_parquet = os.path.join(accel_stage, "cleaned.parquet")
            profile_json = os.path.join(accel_evidence, "profile.json")
            xlsx_path = os.path.join(accel_artifacts, "fin.xlsx")
            docx_path = os.path.join(accel_artifacts, "audit.docx")
            pptx_path = os.path.join(accel_artifacts, "deck.pptx")

            with open(cleaned_csv, "w", encoding="utf-8") as f:
                f.write("id,amount\n10,88.1\n")
            with open(cleaned_parquet, "wb") as f:
                f.write(b"PAR1dataPAR1")
            with open(profile_json, "w", encoding="utf-8") as f:
                f.write("{}")
            with open(xlsx_path, "wb") as f:
                f.write(b"XLSX")
            with open(docx_path, "wb") as f:
                f.write(b"DOCX")
            with open(pptx_path, "wb") as f:
                f.write(b"PPTX")

            accel_outputs = {
                "cleaned_csv": {"path": cleaned_csv, "sha256": ""},
                "cleaned_parquet": {"path": cleaned_parquet, "sha256": ""},
                "profile_json": {"path": profile_json, "sha256": ""},
                "xlsx_fin": {"path": xlsx_path, "sha256": ""},
                "audit_docx": {"path": docx_path, "sha256": ""},
                "deck_pptx": {"path": pptx_path, "sha256": ""},
            }

            with patch("aiwf.flows.cleaning._base_step_start"), patch(
                "aiwf.flows.cleaning._base_artifact_upsert"
            ), patch("aiwf.flows.cleaning._base_step_done"), patch(
                "aiwf.flows.cleaning._base_step_fail"
            ), patch(
                "aiwf.flows.cleaning._try_accel_cleaning",
                return_value={
                    "attempted": True,
                    "ok": True,
                    "response": {
                        "outputs": accel_outputs,
                        "profile": {"rows": 1, "cols": 2},
                    },
                },
            ), patch(
                "aiwf.flows.cleaning._write_fin_xlsx", side_effect=AssertionError("local office writer should not run")
            ), patch(
                "aiwf.flows.cleaning._write_audit_docx", side_effect=AssertionError("local office writer should not run")
            ), patch(
                "aiwf.flows.cleaning._write_deck_pptx", side_effect=AssertionError("local office writer should not run")
            ):
                out = cleaning.run_cleaning(
                    job_id="job-accel-office",
                    actor="test",
                    params=with_job_context(local_job_root, rows=[{"id": 1, "amount": 100.0}]),
                )

            artifacts_by_kind = {artifact["kind"]: artifact["path"] for artifact in out["artifacts"]}
            self.assertEqual(artifacts_by_kind["parquet"], cleaned_parquet)
            self.assertEqual(artifacts_by_kind["xlsx"], xlsx_path)
            self.assertEqual(artifacts_by_kind["docx"], docx_path)
            self.assertEqual(artifacts_by_kind["pptx"], pptx_path)
            self.assertFalse(out["accel"]["used_fallback"])

    def test_materialize_accel_cleaning_artifacts_requires_required_output(self):
        with self.assertRaises(RuntimeError):
            materialize_accel_cleaning_artifacts(
                {},
                params_effective={},
                sha256_file=lambda path: "sha",
            )

    def test_materialize_accel_cleaning_artifacts_skips_missing_optional_outputs(self):
        out = materialize_accel_cleaning_artifacts(
            {
                "cleaned_parquet": {
                    "path": r"D:\tmp\cleaned.parquet",
                    "sha256": "parquet-sha",
                }
            },
            params_effective={},
            sha256_file=lambda path: "sha",
        )

        self.assertEqual(out["cleaned_parquet"], r"D:\tmp\cleaned.parquet")
        self.assertEqual(out["sha_parquet"], "parquet-sha")
        self.assertNotIn("cleaned_csv", out)
        self.assertNotIn("profile_json", out)
        self.assertEqual(
            [artifact["artifact_id"] for artifact in out["core_artifacts"]],
            ["parquet_cleaned_001"],
        )

    def test_run_cleaning_fails_when_local_parquet_invalid_in_strict_mode(self):
        with tempfile.TemporaryDirectory() as tmp:
            local_job_root = os.path.join(tmp, "job")

            def write_bad_parquet(path, rows):
                with open(path, "wb") as f:
                    f.write(b"PARQUET_PLACEHOLDER\n")

            with patch("aiwf.flows.cleaning._base_step_start"), patch(
                "aiwf.flows.cleaning._base_artifact_upsert"
            ), patch("aiwf.flows.cleaning._base_step_done"), patch(
                "aiwf.flows.cleaning._base_step_fail"
            ), patch(
                "aiwf.flows.cleaning._try_accel_cleaning",
                return_value={"attempted": True, "ok": False, "error": "accel unavailable"},
            ), patch(
                "aiwf.flows.cleaning._write_cleaned_parquet", side_effect=write_bad_parquet
            ):
                with self.assertRaises(RuntimeError) as ctx:
                    cleaning.run_cleaning(
                        job_id="job-strict",
                        actor="test",
                        params=with_job_context(local_job_root, local_parquet_strict=True, rows=[{"id": 1, "amount": 100.0}]),
                    )
                self.assertIn("strict mode enabled", str(ctx.exception))

    def test_run_cleaning_fails_on_quality_gate(self):
        with tempfile.TemporaryDirectory() as tmp:
            local_job_root = os.path.join(tmp, "job")

            with patch("aiwf.flows.cleaning._base_step_start"), patch(
                "aiwf.flows.cleaning._base_artifact_upsert"
            ), patch("aiwf.flows.cleaning._base_step_done"), patch(
                "aiwf.flows.cleaning._base_step_fail"
            ), patch(
                "aiwf.flows.cleaning._try_accel_cleaning",
                return_value={"attempted": True, "ok": False, "error": "accel unavailable"},
            ):
                with self.assertRaises(RuntimeError) as ctx:
                    cleaning.run_cleaning(
                        job_id="job-gate",
                        actor="test",
                        params=with_job_context(
                            local_job_root,
                            rows=[{"id": "bad", "amount": "1"}],
                            max_invalid_rows=0,
                        ),
                    )
                self.assertIn("quality gate failed", str(ctx.exception))

    def test_run_cleaning_fails_on_required_missing_ratio_gate(self):
        with tempfile.TemporaryDirectory() as tmp:
            local_job_root = os.path.join(tmp, "job")

            with patch("aiwf.flows.cleaning._base_step_start"), patch(
                "aiwf.flows.cleaning._base_artifact_upsert"
            ), patch("aiwf.flows.cleaning._base_step_done"), patch(
                "aiwf.flows.cleaning._base_step_fail"
            ), patch(
                "aiwf.flows.cleaning._try_accel_cleaning",
                return_value={"attempted": True, "ok": False, "error": "accel unavailable"},
            ):
                with self.assertRaises(RuntimeError) as ctx:
                    cleaning.run_cleaning(
                        job_id="job-required-gate",
                        actor="test",
                        params=with_job_context(
                            local_job_root,
                            rows=[{"name": "alice"}],
                            rules={
                                "platform_mode": "generic",
                            },
                            quality_rules={
                                "required_fields": ["name", "phone"],
                                "max_required_missing_ratio": 0.4,
                            },
                            office_outputs_enabled=False,
                        ),
                    )
                self.assertIn("required_missing_ratio", str(ctx.exception))

    def test_run_cleaning_generic_mode_skips_accel(self):
        with tempfile.TemporaryDirectory() as tmp:
            local_job_root = os.path.join(tmp, "job")

            with patch("aiwf.flows.cleaning._base_step_start"), patch(
                "aiwf.flows.cleaning._base_artifact_upsert"
            ), patch("aiwf.flows.cleaning._base_step_done"), patch(
                "aiwf.flows.cleaning._base_step_fail"
            ), patch(
                "aiwf.flows.cleaning._try_accel_cleaning",
                side_effect=AssertionError("accel should be skipped"),
            ):
                out = cleaning.run_cleaning(
                    job_id="job-generic",
                    actor="test",
                    params=with_job_context(
                        local_job_root,
                        rows=[{"name": "A", "amount": 12}],
                        rules={"platform_mode": "generic"},
                    ),
                )
            self.assertTrue(out["ok"])
            self.assertFalse(out["accel"]["attempted"])

    def test_run_cleaning_with_preprocess_enabled(self):
        with tempfile.TemporaryDirectory() as tmp:
            local_job_root = os.path.join(tmp, "job")
            os.makedirs(local_job_root, exist_ok=True)
            in_csv = os.path.join(local_job_root, "raw.csv")
            with open(in_csv, "w", encoding="utf-8", newline="\n") as f:
                f.write('ID,Amt\n')
                f.write('1,"$10.20"\n')
                f.write('2," 11 "\n')

            with patch("aiwf.flows.cleaning._base_step_start"), patch(
                "aiwf.flows.cleaning._base_artifact_upsert"
            ), patch("aiwf.flows.cleaning._base_step_done"), patch(
                "aiwf.flows.cleaning._base_step_fail"
            ), patch(
                "aiwf.flows.cleaning._try_accel_cleaning",
                return_value={"attempted": True, "ok": False, "error": "accel unavailable"},
            ):
                out = cleaning.run_cleaning(
                    job_id="job-pre",
                    actor="test",
                    params=with_job_context(
                        local_job_root,
                        preprocess={
                            "enabled": True,
                            "input_path": in_csv,
                            "header_map": {"ID": "id", "Amt": "amount"},
                            "amount_fields": ["amount"],
                        },
                    ),
                )
            self.assertTrue(out["ok"])
            self.assertIn("preprocess", out["profile"])
            self.assertEqual(out["profile"]["preprocess"]["summary"]["input_rows"], 2)

    def test_run_cleaning_with_preprocess_pipeline_enabled(self):
        with tempfile.TemporaryDirectory() as tmp:
            local_job_root = os.path.join(tmp, "job")
            os.makedirs(local_job_root, exist_ok=True)
            in_txt = os.path.join(local_job_root, "raw.txt")
            with open(in_txt, "w", encoding="utf-8") as f:
                f.write("ID:1 amount:10\nID:2 amount:11")

            with patch("aiwf.flows.cleaning._base_step_start"), patch(
                "aiwf.flows.cleaning._base_artifact_upsert"
            ), patch("aiwf.flows.cleaning._base_step_done"), patch(
                "aiwf.flows.cleaning._base_step_fail"
            ), patch(
                "aiwf.flows.cleaning._try_accel_cleaning",
                return_value={"attempted": True, "ok": False, "error": "accel unavailable"},
            ):
                out = cleaning.run_cleaning(
                    job_id="job-pre-pipeline",
                    actor="test",
                    params=with_job_context(
                        local_job_root,
                        preprocess={
                            "enabled": True,
                            "input_path": in_txt,
                            "pipeline": {
                                "enabled": True,
                                "stages": [
                                    {"name": "extract", "config": {"input_files": [in_txt]}},
                                    {
                                        "name": "clean",
                                        "config": {
                                            "header_map": {"text": "amount"},
                                            "field_transforms": [{"field": "amount", "op": "extract_regex", "pattern": "amount:(\\d+)", "group": 1}],
                                        },
                                    },
                                ],
                            },
                        },
                    ),
                )
            self.assertTrue(out["ok"])
            self.assertIn("preprocess", out["profile"])
            self.assertEqual(out["profile"]["preprocess"]["mode"], "pipeline")

    def test_clean_rows_use_rust_v2_when_enabled(self):
        fake_resp = Mock()
        fake_resp.status_code = 200
        fake_resp.json.return_value = {
            "ok": True,
            "rows": [{"id": 1, "amount": 10.0}],
            "quality": {"input_rows": 1, "output_rows": 1, "invalid_rows": 0, "filtered_rows": 0, "duplicate_rows_removed": 0},
            "trace_id": "t1",
        }
        with patch("requests.post", return_value=fake_resp):
            out = cleaning._clean_rows(
                [{"id": "1", "amount": "10"}],
                {"rules": {"use_rust_v2": True}},
            )
        self.assertEqual(out["rows"], [{"id": 1, "amount": 10.0}])
        self.assertTrue(out["quality"].get("rust_v2_used"))
        self.assertEqual(out["quality"].get("rust_v2_trace_id"), "t1")
        self.assertEqual(out["quality"].get("cleaning_spec_version"), "cleaning_spec.v2")
        self.assertEqual(out["execution_mode"], "rust_v2")
        self.assertEqual(out["eligibility_reason"], "eligible")
        self.assertEqual(out["execution_audit"]["operator"], "transform_rows_v3")

    def test_clean_rows_use_rust_v2_for_generic_rules_after_spec_compile(self):
        fake_resp = Mock()
        fake_resp.status_code = 200
        fake_resp.json.return_value = {
            "ok": True,
            "rows": [{"customer_name": "alice", "amount": 10.0}],
            "quality": {"input_rows": 1, "output_rows": 1, "invalid_rows": 0, "filtered_rows": 0, "duplicate_rows_removed": 0},
            "trace_id": "tg1",
            "audit": {"schema": "transform_rows_v2.audit.v1"},
        }
        with patch("requests.post", return_value=fake_resp):
            out = cleaning._clean_rows(
                [{"CustomerName": "Alice", "Amount": "10"}],
                {
                    "rules": {
                        "use_rust_v2": True,
                        "platform_mode": "generic",
                        "rename_map": {"CustomerName": "customer_name", "Amount": "amount"},
                        "casts": {"amount": "float"},
                    }
                },
        )
        self.assertEqual(out["rows"], [{"customer_name": "alice", "amount": 10.0}])
        self.assertTrue(out["quality"].get("rust_v2_used"))
        self.assertEqual(out["execution_audit"]["schema"], "transform_rows_v2.audit.v1")

    def test_clean_rows_fallback_when_rust_v2_unavailable(self):
        with patch("requests.post", side_effect=RuntimeError("unreachable")):
            out = cleaning._clean_rows(
                [{"id": "1", "amount": "10"}],
                {"rules": {"use_rust_v2": True}},
            )
        self.assertEqual(out["rows"], [{"id": 1, "amount": 10.0}])
        self.assertFalse(out["quality"].get("rust_v2_used"))
        self.assertEqual(out["execution_mode"], "python_legacy")
        self.assertEqual(out["eligibility_reason"], "rust_v2_error")
        self.assertEqual(out["execution_audit"]["schema"], "python_cleaning.audit.v1")
        self.assertIn("rust_v2_error", out["execution_audit"])

    def test_clean_rows_marks_python_legacy_when_flag_disabled(self):
        out = cleaning._clean_rows(
            [{"id": "1", "amount": "10"}],
            {"rules": {}},
        )
        self.assertEqual(out["execution_mode"], "python_legacy")
        self.assertEqual(out["eligibility_reason"], "mode_off")
        self.assertEqual(out["requested_rust_v2_mode"], "off")
        self.assertEqual(out["effective_rust_v2_mode"], "off")
        self.assertFalse(out["verify_on_default"])
        self.assertEqual(out["shadow_compare"]["status"], "skipped")
        self.assertEqual(out["shadow_compare"]["skipped_reason"], "mode_off")

    def test_python_reason_counts_maps_invalid_id_amount_to_required_missing(self):
        counts = cleaning._python_reason_counts(
            {
                "rule_hits": {
                    "cast_failed": 1,
                    "invalid_id": 2,
                    "invalid_amount": 3,
                    "required_failed": 4,
                }
            }
        )
        self.assertEqual(counts["cast_failed"], 1)
        self.assertEqual(counts["required_missing"], 9)

    def test_clean_rows_shadow_mode_reports_match(self):
        fake_resp = Mock()
        fake_resp.status_code = 200
        fake_resp.json.return_value = {
            "ok": True,
            "rows": [{"id": 1, "amount": 10.0}],
            "quality": {
                "input_rows": 1,
                "output_rows": 1,
                "invalid_rows": 0,
                "filtered_rows": 0,
                "duplicate_rows_removed": 0,
                "required_missing_ratio": 0.0,
            },
            "trace_id": "shadow-ok",
            "audit": {
                "schema": "transform_rows_v2.audit.v1",
                "reason_counts": {
                    "invalid_object": 0,
                    "cast_failed": 0,
                    "required_missing": 0,
                    "filter_rejected": 0,
                    "duplicate_removed": 0,
                },
                "reason_samples": {
                    "invalid_object": [],
                    "cast_failed": [],
                    "required_missing": [],
                    "filter_rejected": [],
                    "duplicate_removed": [],
                },
                "limits": {"sample_limit": 5},
            },
        }
        with patch.dict(os.environ, {"AIWF_CLEANING_RUST_V2_MODE": "shadow"}, clear=False), patch("requests.post", return_value=fake_resp):
            out = cleaning._clean_rows(
                [{"id": "1", "amount": "10"}],
                {"rules": {}},
            )
        self.assertEqual(out["execution_mode"], "python_legacy")
        self.assertEqual(out["requested_rust_v2_mode"], "shadow")
        self.assertEqual(out["effective_rust_v2_mode"], "shadow")
        self.assertFalse(out["verify_on_default"])
        self.assertEqual(out["shadow_compare"]["status"], "matched")
        self.assertTrue(out["shadow_compare"]["matched"])

    def test_clean_rows_shadow_mode_reports_mismatch(self):
        fake_resp = Mock()
        fake_resp.status_code = 200
        fake_resp.json.return_value = {
            "ok": True,
            "rows": [{"id": 1, "amount": 999.0}],
            "quality": {
                "input_rows": 1,
                "output_rows": 1,
                "invalid_rows": 0,
                "filtered_rows": 0,
                "duplicate_rows_removed": 0,
                "required_missing_ratio": 0.0,
            },
            "trace_id": "shadow-bad",
            "audit": {
                "schema": "transform_rows_v2.audit.v1",
                "reason_counts": {
                    "invalid_object": 0,
                    "cast_failed": 0,
                    "required_missing": 0,
                    "filter_rejected": 0,
                    "duplicate_removed": 0,
                },
                "reason_samples": {
                    "invalid_object": [],
                    "cast_failed": [],
                    "required_missing": [],
                    "filter_rejected": [],
                    "duplicate_removed": [],
                },
                "limits": {"sample_limit": 5},
            },
        }
        with patch.dict(os.environ, {"AIWF_CLEANING_RUST_V2_MODE": "shadow"}, clear=False), patch("requests.post", return_value=fake_resp):
            out = cleaning._clean_rows(
                [{"id": "1", "amount": "10"}],
                {"rules": {}},
            )
        self.assertEqual(out["execution_mode"], "python_legacy")
        self.assertEqual(out["shadow_compare"]["status"], "mismatched")
        self.assertGreater(out["shadow_compare"]["mismatch_count"], 0)

    def test_clean_rows_default_mode_prefers_rust(self):
        fake_resp = Mock()
        fake_resp.status_code = 200
        fake_resp.json.return_value = {
            "ok": True,
            "rows": [{"id": 1, "amount": 10.0}],
            "quality": {
                "input_rows": 1,
                "output_rows": 1,
                "invalid_rows": 0,
                "filtered_rows": 0,
                "duplicate_rows_removed": 0,
                "required_missing_ratio": 0.0,
            },
            "trace_id": "default-ok",
            "audit": {"schema": "transform_rows_v2.audit.v1"},
        }
        with patch.dict(os.environ, {"AIWF_CLEANING_RUST_V2_MODE": "default"}, clear=False), patch("requests.post", return_value=fake_resp):
            out = cleaning._clean_rows(
                [{"id": "1", "amount": "10"}],
                {"rules": {}},
            )
        self.assertEqual(out["execution_mode"], "rust_v2")
        self.assertEqual(out["requested_rust_v2_mode"], "default")
        self.assertEqual(out["effective_rust_v2_mode"], "default")
        self.assertFalse(out["verify_on_default"])
        self.assertEqual(out["shadow_compare"]["status"], "skipped")
        self.assertEqual(out["shadow_compare"]["skipped_reason"], "default_without_verify")

    def test_clean_rows_default_mode_with_verify_reports_match(self):
        fake_resp = Mock()
        fake_resp.status_code = 200
        fake_resp.json.return_value = {
            "ok": True,
            "rows": [{"id": 1, "amount": 10.0}],
            "quality": {
                "input_rows": 1,
                "output_rows": 1,
                "invalid_rows": 0,
                "filtered_rows": 0,
                "duplicate_rows_removed": 0,
                "required_missing_ratio": 0.0,
            },
            "trace_id": "default-verify-ok",
            "audit": {
                "schema": "transform_rows_v2.audit.v1",
                "reason_counts": {
                    "invalid_object": 0,
                    "cast_failed": 0,
                    "required_missing": 0,
                    "filter_rejected": 0,
                    "duplicate_removed": 0,
                },
                "reason_samples": {
                    "invalid_object": [],
                    "cast_failed": [],
                    "required_missing": [],
                    "filter_rejected": [],
                    "duplicate_removed": [],
                },
                "limits": {"sample_limit": 5},
            },
        }
        with patch.dict(
            os.environ,
            {
                "AIWF_CLEANING_RUST_V2_MODE": "default",
                "AIWF_CLEANING_RUST_V2_VERIFY_ON_DEFAULT": "true",
            },
            clear=False,
        ), patch("requests.post", return_value=fake_resp):
            out = cleaning._clean_rows(
                [{"id": "1", "amount": "10"}],
                {"rules": {}},
            )
        self.assertEqual(out["execution_mode"], "rust_v2")
        self.assertEqual(out["requested_rust_v2_mode"], "default")
        self.assertEqual(out["effective_rust_v2_mode"], "default")
        self.assertTrue(out["verify_on_default"])
        self.assertEqual(out["shadow_compare"]["status"], "matched")
        self.assertTrue(out["shadow_compare"]["matched"])

    def test_clean_rows_default_mode_with_verify_keeps_rust_on_mismatch(self):
        fake_resp = Mock()
        fake_resp.status_code = 200
        fake_resp.json.return_value = {
            "ok": True,
            "rows": [{"id": 1, "amount": 999.0}],
            "quality": {
                "input_rows": 1,
                "output_rows": 1,
                "invalid_rows": 0,
                "filtered_rows": 0,
                "duplicate_rows_removed": 0,
                "required_missing_ratio": 0.0,
            },
            "trace_id": "default-verify-bad",
            "audit": {
                "schema": "transform_rows_v2.audit.v1",
                "reason_counts": {
                    "invalid_object": 0,
                    "cast_failed": 0,
                    "required_missing": 0,
                    "filter_rejected": 0,
                    "duplicate_removed": 0,
                },
                "reason_samples": {
                    "invalid_object": [],
                    "cast_failed": [],
                    "required_missing": [],
                    "filter_rejected": [],
                    "duplicate_removed": [],
                },
                "limits": {"sample_limit": 5},
            },
        }
        with patch.dict(
            os.environ,
            {
                "AIWF_CLEANING_RUST_V2_MODE": "default",
                "AIWF_CLEANING_RUST_V2_VERIFY_ON_DEFAULT": "true",
            },
            clear=False,
        ), patch("requests.post", return_value=fake_resp):
            out = cleaning._clean_rows(
                [{"id": "1", "amount": "10"}],
                {"rules": {}},
            )
        self.assertEqual(out["execution_mode"], "python_legacy")
        self.assertEqual(out["eligibility_reason"], "shadow_compare_mismatch")
        self.assertEqual(out["requested_rust_v2_mode"], "default")
        self.assertEqual(out["effective_rust_v2_mode"], "default")
        self.assertTrue(out["verify_on_default"])
        self.assertEqual(out["shadow_compare"]["status"], "mismatched")
        self.assertGreater(out["shadow_compare"]["mismatch_count"], 0)

    def test_clean_rows_local_standalone_defaults_to_default_verify(self):
        fake_resp = Mock()
        fake_resp.status_code = 200
        fake_resp.json.return_value = {
            "ok": True,
            "rows": [{"id": 1, "amount": 10.0}],
            "quality": {
                "input_rows": 1,
                "output_rows": 1,
                "invalid_rows": 0,
                "filtered_rows": 0,
                "duplicate_rows_removed": 0,
                "required_missing_ratio": 0.0,
            },
            "trace_id": "standalone-default-ok",
            "audit": {
                "schema": "transform_rows_v2.audit.v1",
                "reason_counts": {
                    "invalid_object": 0,
                    "cast_failed": 0,
                    "required_missing": 0,
                    "filter_rejected": 0,
                    "duplicate_removed": 0,
                },
                "reason_samples": {
                    "invalid_object": [],
                    "cast_failed": [],
                    "required_missing": [],
                    "filter_rejected": [],
                    "duplicate_removed": [],
                },
                "limits": {"sample_limit": 5},
            },
        }
        with patch("requests.post", return_value=fake_resp):
            out = cleaning._clean_rows(
                [{"id": "1", "amount": "10"}],
                {"local_standalone": True},
            )
        self.assertEqual(out["execution_mode"], "rust_v2")
        self.assertEqual(out["requested_rust_v2_mode"], "default")
        self.assertEqual(out["effective_rust_v2_mode"], "default")
        self.assertTrue(out["verify_on_default"])
        self.assertEqual(out["shadow_compare"]["status"], "matched")

    def test_rules_use_rust_v2_false_overrides_default_mode(self):
        with patch.dict(os.environ, {"AIWF_CLEANING_RUST_V2_MODE": "default"}, clear=False), patch("requests.post") as post:
            out = cleaning._clean_rows(
                [{"id": "1", "amount": "10"}],
                {"rules": {"use_rust_v2": False}},
            )
        post.assert_not_called()
        self.assertEqual(out["execution_mode"], "python_legacy")
        self.assertEqual(out["eligibility_reason"], "forced_python")
        self.assertEqual(out["requested_rust_v2_mode"], "default")
        self.assertEqual(out["effective_rust_v2_mode"], "force_python")
        self.assertFalse(out["verify_on_default"])
        self.assertEqual(out["shadow_compare"]["status"], "skipped")
        self.assertEqual(out["shadow_compare"]["skipped_reason"], "forced_python")

    def test_clean_rows_default_mode_falls_back_on_rust_error(self):
        with patch.dict(os.environ, {"AIWF_CLEANING_RUST_V2_MODE": "default"}, clear=False), patch("requests.post", side_effect=RuntimeError("unreachable")):
            out = cleaning._clean_rows(
                [{"id": "1", "amount": "10"}],
                {"rules": {}},
            )
        self.assertEqual(out["execution_mode"], "python_legacy")
        self.assertEqual(out["eligibility_reason"], "rust_v2_error")
        self.assertEqual(out["shadow_compare"]["status"], "rust_error")

    def test_run_cleaning_surfaces_execution_report_on_local_path(self):
        with tempfile.TemporaryDirectory() as tmp:
            local_job_root = os.path.join(tmp, "job")

            def write_valid_parquet(path, rows):
                with open(path, "wb") as f:
                    f.write(b"PAR1dataPAR1")

            with patch("aiwf.flows.cleaning._base_step_start"), patch(
                "aiwf.flows.cleaning._base_artifact_upsert"
            ), patch("aiwf.flows.cleaning._base_step_done"), patch(
                "aiwf.flows.cleaning._base_step_fail"
            ), patch(
                "aiwf.flows.cleaning._try_accel_cleaning",
                return_value={"attempted": True, "ok": False, "error": "accel unavailable"},
            ), patch(
                "aiwf.flows.cleaning._write_cleaned_parquet", side_effect=write_valid_parquet
            ):
                out = cleaning.run_cleaning(
                    job_id="job-execution-report",
                    actor="test",
                    params=with_job_context(local_job_root, rows=[{"id": 1, "amount": 10.0}], rules={"use_rust_v2": False}),
                )

            self.assertIn("execution", out)
            self.assertEqual(out["execution"]["execution_mode"], "python_legacy")
            self.assertEqual(out["profile"]["execution"]["eligibility_reason"], "forced_python")
            self.assertEqual(out["execution"]["requested_rust_v2_mode"], "off")
            self.assertEqual(out["execution"]["effective_rust_v2_mode"], "force_python")
            self.assertFalse(out["execution"]["verify_on_default"])
            self.assertIn("shadow_compare", out["execution"])

    def test_run_cleaning_local_standalone_skips_base_transport(self):
        with tempfile.TemporaryDirectory() as tmp:
            local_job_root = os.path.join(tmp, "job")

            def write_valid_parquet(path, rows):
                with open(path, "wb") as f:
                    f.write(b"PAR1dataPAR1")

            fake_resp = Mock()
            fake_resp.status_code = 200
            fake_resp.json.return_value = {
                "ok": True,
                "rows": [{"id": 1, "amount": 10.0}],
                "quality": {
                    "input_rows": 1,
                    "output_rows": 1,
                    "invalid_rows": 0,
                    "filtered_rows": 0,
                    "duplicate_rows_removed": 0,
                    "required_missing_ratio": 0.0,
                },
                "trace_id": "standalone-run-ok",
                "audit": {
                    "schema": "transform_rows_v2.audit.v1",
                    "reason_counts": {
                        "invalid_object": 0,
                        "cast_failed": 0,
                        "required_missing": 0,
                        "filter_rejected": 0,
                        "duplicate_removed": 0,
                    },
                    "reason_samples": {
                        "invalid_object": [],
                        "cast_failed": [],
                        "required_missing": [],
                        "filter_rejected": [],
                        "duplicate_removed": [],
                    },
                    "limits": {"sample_limit": 5},
                },
            }

            with patch("aiwf.flows.cleaning._base_step_start") as step_start, patch(
                "aiwf.flows.cleaning._base_artifact_upsert"
            ) as artifact_upsert, patch("aiwf.flows.cleaning._base_step_done") as step_done, patch(
                "aiwf.flows.cleaning._base_step_fail"
            ) as step_fail, patch(
                "aiwf.flows.cleaning._try_accel_cleaning",
                return_value={"attempted": True, "ok": False, "error": "accel unavailable"},
            ), patch(
                "aiwf.flows.cleaning._write_cleaned_parquet", side_effect=write_valid_parquet
            ), patch("requests.post", return_value=fake_resp):
                out = cleaning.run_cleaning(
                    job_id="job-standalone",
                    actor="test",
                    params=with_job_context(local_job_root, rows=[{"id": "1", "amount": "10"}], local_standalone=True),
                )

            step_start.assert_not_called()
            artifact_upsert.assert_not_called()
            step_done.assert_not_called()
            step_fail.assert_not_called()
            self.assertEqual(out["execution"]["execution_mode"], "rust_v2")
            self.assertEqual(out["execution"]["requested_rust_v2_mode"], "default")
            self.assertTrue(out["execution"]["verify_on_default"])

    def test_write_cleaned_csv_quotes_delimited_fields(self):
        with tempfile.TemporaryDirectory() as tmp:
            csv_path = os.path.join(tmp, "cleaned.csv")
            cleaning._write_cleaned_csv(csv_path, [{"id": 1, "note": "a,b"}])
            with open(csv_path, "r", encoding="utf-8") as f:
                text = f.read()
            self.assertIn('"a,b"', text)


if __name__ == "__main__":
    unittest.main()
